from __future__ import annotations

import math
import hashlib
import json
import logging
import shutil
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from sen_cinematic_engine_monster import MonsterDeck, SLIDES as MONSTER_SLIDES

log = logging.getLogger("sen_deck")
logging.basicConfig(level=logging.INFO, format="%(message)s")

# ════════════════════════════════════════════════════════════════
#  PATHS
# ════════════════════════════════════════════════════════════════
ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "SEN Final.pptx"
ASSETS = ROOT / "sen_final_assets"
DOC_MEDIA = ROOT / "_sen_doc_media"
DERIVED = ASSETS / "_derived"

def first_existing_path(candidates: Iterable[Path], *, fallback: Path) -> Path:
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return fallback


SEN_ROOT = first_existing_path(
    [
        Path(r"D:\sen"),
        ROOT / "sen_project",
        ROOT / "sen_assets",
    ],
    fallback=ROOT,
)


def resolve_asset(*candidates: Path | str) -> Path:
    resolved: list[Path] = []
    for candidate in candidates:
        path = candidate if isinstance(candidate, Path) else Path(candidate)
        if path.is_absolute():
            resolved.append(path)
        else:
            resolved.extend(
                [
                    ROOT / path,
                    SEN_ROOT / path,
                    ROOT / "sen_pro_assets" / path.name,
                    ROOT / "sen_presentation_assets" / path.name,
                    ROOT / "sen_final_assets" / path.name,
                ]
            )
    return first_existing_path(resolved, fallback=resolved[0])


LOGO_MARK = resolve_asset("assets/logos/logo.png", "logo.png")
LOGO_WORD = resolve_asset("assets/logos/ourLogo.png", "ourLogo.png")
SHOT_HOME = resolve_asset("sen-home-screen.png", "home.png")
SHOT_LOGIN = resolve_asset("sen-login.png", "login.png")
SHOT_SECTIONS = resolve_asset("sen-sections.png", "sections.png")
SHOT_ACCOUNT = resolve_asset("sen-account-type.png", "account-type.png", "11-accounts.png")

UNIV_LOGO = DOC_MEDIA / "image1.jpg"
AGILE_DIAGRAM = DOC_MEDIA / "image3.png"
TIMELINE_CHART = DOC_MEDIA / "image4.png"
AUTH_SEQUENCE = DOC_MEDIA / "image5.png"
USECASE_USER = DOC_MEDIA / "image21.png"
USECASE_WORKER = DOC_MEDIA / "image22.png"
USECASE_STORE = DOC_MEDIA / "image23.png"
ERD_DIAGRAM = DOC_MEDIA / "image24.png"

# ════════════════════════════════════════════════════════════════
#  DESIGN SYSTEM
# ════════════════════════════════════════════════════════════════
SLIDE_W, SLIDE_H = 13.333, 7.5
PX_W, PX_H = 1920, 1080
DISPLAY = "Segoe UI"
BODY = "Segoe UI"
TOTAL_SLIDES = 34

SLIDE_CHAPTERS = [
    "الافتتاح", "الافتتاح", "الملخص", "المشكلة", "المشكلة", "تعريف SEN", "تعريف SEN",
    "النطاق", "لماذا الآن؟", "المنهجية", "أصحاب المصلحة", "التحليل", "النموذج",
    "المنتج", "المنتج", "المتطلبات", "المتطلبات", "المتطلبات", "المتطلبات",
    "المتطلبات", "الجدوى", "الجدوى", "المخططات", "المخططات", "المخططات", "الثقة",
    "الأعمال", "الأعمال", "التميّز", "التطوير", "القيمة", "الدفاع", "الختام", "الختام",
]

CHAPTER_ACCENTS = {
    "الافتتاح": "aqua",
    "الملخص": "royal",
    "المشكلة": "coral",
    "تعريف SEN": "violet",
    "النطاق": "sky",
    "لماذا الآن؟": "amber",
    "المنهجية": "mint",
    "أصحاب المصلحة": "violet",
    "التحليل": "coral",
    "النموذج": "royal",
    "المنتج": "aqua",
    "المتطلبات": "electric",
    "الجدوى": "amber",
    "المخططات": "sky",
    "الثقة": "mint",
    "الأعمال": "amber",
    "التميّز": "violet",
    "التطوير": "emerald",
    "القيمة": "aqua",
    "الدفاع": "royal",
    "الختام": "aqua",
}

COLORS = {
    "void": (3, 5, 12), "night": (8, 13, 24), "navy": (13, 27, 53),
    "deep": (18, 35, 66), "obsidian": (12, 14, 22),
    "royal": (41, 100, 255), "electric": (76, 158, 255),
    "sky": (111, 186, 255), "ice": (180, 220, 255),
    "aqua": (98, 242, 255), "mint": (122, 255, 214),
    "emerald": (46, 213, 150), "neon": (0, 255, 180),
    "violet": (148, 120, 255), "purple": (100, 70, 200),
    "magenta": (255, 80, 200), "lavender": (180, 160, 255),
    "coral": (255, 117, 96), "amber": (255, 189, 87),
    "gold": (255, 215, 0), "rose": (255, 136, 160),
    "paper": (247, 249, 253), "snow": (252, 253, 255),
    "mist": (216, 225, 239), "ash": (160, 170, 185),
    "grid": (112, 129, 170), "ink": (18, 26, 40), "slate": (44, 56, 80),
}


def rgb(name: str) -> RGBColor:
    return RGBColor(*COLORS[name])


def ensure_assets() -> None:
    ASSETS.mkdir(exist_ok=True)
    DERIVED.mkdir(exist_ok=True)


def chapter_for_slide(num: int) -> str:
    if 1 <= num <= len(SLIDE_CHAPTERS):
        return SLIDE_CHAPTERS[num - 1]
    return "SEN"


def accent_for_slide(num: int) -> str:
    return CHAPTER_ACCENTS.get(chapter_for_slide(num), "aqua")


def mood_for_slide(num: int) -> str:
    light_chapters = {"الملخص", "النطاق", "أصحاب المصلحة", "المخططات", "الثقة", "التطوير", "الدفاع"}
    return "light" if chapter_for_slide(num) in light_chapters else "dark"


def crop_asset(path: Path, box: tuple[int, int, int, int], label: str) -> Path:
    if not path.exists():
        log.warning("asset missing for crop: %s", path)
        return path
    signature = hashlib.sha1(f"{path}:{box}:{label}".encode("utf-8")).hexdigest()[:12]
    out = DERIVED / f"{label}-{signature}.png"
    if out.exists():
        return out
    with Image.open(path) as img:
        cropped = img.crop(box)
        cropped.save(out)
    return out


# ════════════════════════════════════════════════════════════════
#  BACKGROUND ENGINE — CORE
# ════════════════════════════════════════════════════════════════
def _layer() -> tuple[Image.Image, ImageDraw.Draw]:
    img = Image.new("RGBA", (PX_W, PX_H), (0, 0, 0, 0))
    return img, ImageDraw.Draw(img)


def _composite(base: Image.Image, layer: Image.Image, blur: int = 0) -> Image.Image:
    if blur > 0:
        layer = layer.filter(ImageFilter.GaussianBlur(blur))
    return Image.alpha_composite(base.convert("RGBA"), layer)


def gradient(top, bottom) -> Image.Image:
    img = Image.new("RGB", (PX_W, PX_H), top)
    draw = ImageDraw.Draw(img)
    for y in range(PX_H):
        t = y / max(PX_H - 1, 1)
        c = tuple(int(top[i] * (1 - t) + bottom[i] * t) for i in range(3))
        draw.line((0, y, PX_W, y), fill=c)
    return img


def gradient_radial(center_color, edge_color, *, cx=960, cy=540) -> Image.Image:
    img = Image.new("RGB", (PX_W, PX_H), edge_color)
    draw = ImageDraw.Draw(img)
    max_r = int(math.hypot(PX_W, PX_H))
    for r in range(max_r, 0, -2):
        t = r / max_r
        c = tuple(int(center_color[i] * (1 - t) + edge_color[i] * t) for i in range(3))
        draw.ellipse((cx - r, cy - r, cx + r, cy + r), fill=c)
    return img


def gradient_three(top, mid, bottom) -> Image.Image:
    img = Image.new("RGB", (PX_W, PX_H), top)
    draw = ImageDraw.Draw(img)
    half = PX_H // 2
    for y in range(PX_H):
        if y < half:
            t = y / max(half - 1, 1)
            c = tuple(int(top[i] * (1 - t) + mid[i] * t) for i in range(3))
        else:
            t = (y - half) / max(half - 1, 1)
            c = tuple(int(mid[i] * (1 - t) + bottom[i] * t) for i in range(3))
        draw.line((0, y, PX_W, y), fill=c)
    return img


# ════════════════════════════════════════════════════════════════
#  BACKGROUND ENGINE — EFFECTS (25 effects)
# ════════════════════════════════════════════════════════════════
def add_glow(img, bbox, color, alpha, blur, *, ellipse=True):
    layer, draw = _layer()
    if ellipse:
        draw.ellipse(bbox, fill=(*COLORS[color], alpha))
    else:
        draw.rounded_rectangle(bbox, radius=44, fill=(*COLORS[color], alpha))
    return _composite(img, layer, blur)


def add_grid(img, step, alpha, *, color="grid"):
    layer, draw = _layer()
    c = (*COLORS[color], alpha)
    for x in range(0, PX_W, step):
        draw.line((x, 0, x, PX_H), fill=c, width=1)
    for y in range(0, PX_H, step):
        draw.line((0, y, PX_W, y), fill=c, width=1)
    return _composite(img, layer)


def add_diagonal_lines(img, alpha=28, *, color="sky"):
    layer, draw = _layer()
    for i in range(-8, 24):
        x = i * 124
        draw.line((x, PX_H, x + 780, 0), fill=(*COLORS[color], alpha), width=2)
    return _composite(img, layer, 2)


def add_arcs(img, alpha=28, *, color="aqua", x_shift=0):
    layer, draw = _layer()
    for i in range(9):
        draw.arc((900 + x_shift - i * 58, 90 + i * 34, 1950 + x_shift, 1020 + i * 16), start=182, end=334, fill=(*COLORS[color], alpha), width=3)
    return _composite(img, layer, 1)


def add_blueprint(img, *, color="sky"):
    layer, draw = _layer()
    for i in range(13):
        x = 120 + i * 128
        draw.line((x, 130, x + 280, 960), fill=(*COLORS[color], 20), width=2)
    for i in range(17):
        y = 150 + i * 42
        draw.line((940, y, 1840, y), fill=(*COLORS["mist"], 14), width=1)
    return _composite(img, layer, 1)


def add_nodes(img, *, color="aqua"):
    layer, draw = _layer()
    dots = [(1540, 210), (1380, 350), (1240, 520), (1020, 708), (760, 560), (520, 720)]
    for x, y in dots:
        draw.ellipse((x - 14, y - 14, x + 14, y + 14), fill=(*COLORS[color], 190))
        draw.ellipse((x - 22, y - 22, x + 22, y + 22), outline=(*COLORS[color], 60), width=2)
    for (x1, y1), (x2, y2) in zip(dots[:-1], dots[1:]):
        draw.line((x1, y1, x2, y2), fill=(*COLORS["sky"], 78), width=4)
    return _composite(img, layer, 1)


def add_particles(img, *, color="paper", density=1800, alpha_max=26):
    layer, draw = _layer()
    for i in range(density):
        x = (i * 37 + i * i * 11) % PX_W
        y = (i * 89 + i * i * 7) % PX_H
        a = 3 + (i * 17) % max(alpha_max, 4)
        r = 1 + (i % 4 == 0) + (i % 7 == 0)
        draw.ellipse((x, y, x + r, y + r), fill=(*COLORS[color], a))
    return _composite(img, layer)


def add_waves(img, *, color="aqua", alpha=32, amplitude=34, frequency=2.8, baseline=760):
    layer, draw = _layer()
    for band in range(5):
        pts = []
        for x in range(0, PX_W + 1, 14):
            y = baseline + band * 22 + int(math.sin((x / PX_W) * math.pi * frequency + band * 0.9) * amplitude)
            pts.append((x, y))
        draw.line(pts, fill=(*COLORS[color], alpha - band * 4), width=3)
    return _composite(img, layer, 1)


def add_rings(img, *, color="royal", alpha=28, center=(1540, 260), start_radius=120, count=6):
    layer, draw = _layer()
    cx, cy = center
    for i in range(count):
        r = start_radius + i * 58
        draw.ellipse((cx - r, cy - r, cx + r, cy + r), outline=(*COLORS[color], alpha - i * 3), width=3)
    return _composite(img, layer, 1)


def add_columns(img, *, color="violet", alpha=24, start=1040):
    layer, draw = _layer()
    for i in range(7):
        x = start + i * 100
        draw.rounded_rectangle((x, 110 + i * 8, x + 8, 900 - i * 26), radius=4, fill=(*COLORS[color], alpha + i))
    return _composite(img, layer, 6)


def add_mesh(img, *, color="aqua", alpha=18, x_offset=0):
    layer, draw = _layer()
    for row in range(7):
        pts = []
        for col in range(-1, 10):
            x = 920 + x_offset + col * 140
            y = 140 + row * 110 + int(math.sin((col + row * 0.5) * 0.7) * 16)
            pts.append((x, y))
        draw.line(pts, fill=(*COLORS[color], alpha), width=2)
    for col in range(8):
        pts = []
        for row in range(-1, 9):
            x = 940 + x_offset + col * 140 + int(math.sin((row + col * 0.4) * 0.8) * 16)
            y = 110 + row * 110
            pts.append((x, y))
        draw.line(pts, fill=(*COLORS["sky"], alpha), width=2)
    return _composite(img, layer, 1)


def add_noise(img, amount=10):
    layer, draw = _layer()
    for i in range(4200):
        x, y = (i * 41) % PX_W, (i * 79) % PX_H
        draw.point((x, y), fill=(255, 255, 255, (i * 17) % amount))
    return _composite(img, layer)


# ── NEW CINEMATIC EFFECTS ─────────────────────────────────────

def add_bokeh(img, *, colors=("aqua", "sky", "violet", "electric"), count=40, alpha_range=(12, 40)):
    layer, draw = _layer()
    for i in range(count):
        x = (i * 197 + i * i * 13) % PX_W
        y = (i * 131 + i * i * 9) % PX_H
        r = 18 + (i * 23) % 70
        c = COLORS[colors[i % len(colors)]]
        a = alpha_range[0] + (i * 19) % (alpha_range[1] - alpha_range[0])
        draw.ellipse((x - r, y - r, x + r, y + r), fill=(*c, a))
    return _composite(img, layer, 18)


def add_aurora(img, *, colors=("aqua", "mint", "sky", "violet"), bands=5, amplitude=80, blur_r=45):
    layer, draw = _layer()
    band_h = PX_H // (bands + 2)
    for b in range(bands):
        c = COLORS[colors[b % len(colors)]]
        base_y = 120 + b * band_h
        pts = []
        for x in range(0, PX_W + 1, 8):
            y = base_y + int(math.sin(x / 220.0 + b * 1.3) * amplitude + math.cos(x / 340.0 + b * 0.7) * amplitude * 0.5)
            pts.append((x, y))
        for w in range(40):
            shifted = [(px, py + w) for px, py in pts]
            draw.line(shifted, fill=(*c, 18 - w // 3), width=2)
    return _composite(img, layer, blur_r)


def add_constellation(img, *, color="paper", count=45, connect_dist=210, alpha=22):
    layer, draw = _layer()
    stars = []
    for i in range(count):
        x = (i * 211 + i * i * 17) % PX_W
        y = (i * 157 + i * i * 11) % PX_H
        stars.append((x, y))
        r = 2 + (i % 3)
        draw.ellipse((x - r, y - r, x + r, y + r), fill=(*COLORS[color], alpha + 20))
    for i, (x1, y1) in enumerate(stars):
        for x2, y2 in stars[i + 1:]:
            d = math.hypot(x2 - x1, y2 - y1)
            if d < connect_dist:
                a = int(alpha * (1 - d / connect_dist))
                draw.line((x1, y1, x2, y2), fill=(*COLORS[color], max(a, 3)), width=1)
    return _composite(img, layer, 1)


def add_hexgrid(img, *, color="sky", size=58, alpha=14):
    layer, draw = _layer()
    c = (*COLORS[color], alpha)
    h = size * math.sqrt(3)
    for row in range(-1, int(PX_H / h) + 2):
        for col in range(-1, int(PX_W / (size * 1.5)) + 2):
            cx = col * size * 1.5
            cy = row * h + (h / 2 if col % 2 else 0)
            pts = []
            for k in range(6):
                angle = math.radians(60 * k + 30)
                pts.append((cx + size * 0.55 * math.cos(angle), cy + size * 0.55 * math.sin(angle)))
            draw.polygon(pts, outline=c)
    return _composite(img, layer, 1)


def add_topographic(img, *, color="electric", lines=14, alpha=16):
    layer, draw = _layer()
    for i in range(lines):
        pts = []
        base_y = 80 + i * (PX_H // (lines + 1))
        for x in range(0, PX_W + 1, 10):
            y = base_y + int(math.sin(x / 180.0 + i * 0.8) * 30 + math.sin(x / 90.0 + i * 1.6) * 15)
            pts.append((x, y))
        draw.line(pts, fill=(*COLORS[color], alpha - i % 4), width=2)
    return _composite(img, layer, 1)


def add_lens_flare(img, cx=1500, cy=300, *, color="paper", alpha=55, size=180):
    layer, draw = _layer()
    for r in range(size, 0, -3):
        a = int(alpha * (1 - r / size) ** 1.5)
        draw.ellipse((cx - r, cy - r, cx + r, cy + r), fill=(*COLORS[color], a))
    ghosts = [(cx - 300, cy + 200, 40), (cx - 500, cy + 350, 25), (cx - 180, cy + 100, 55)]
    for gx, gy, gr in ghosts:
        draw.ellipse((gx - gr, gy - gr, gx + gr, gy + gr), fill=(*COLORS["sky"], 18))
    return _composite(img, layer, 8)


def add_radial_burst(img, cx=960, cy=540, *, color="royal", rays=28, alpha=14, length=600):
    layer, draw = _layer()
    for i in range(rays):
        angle = math.radians(i * (360 / rays))
        x2 = cx + int(math.cos(angle) * length)
        y2 = cy + int(math.sin(angle) * length)
        draw.line((cx, cy, x2, y2), fill=(*COLORS[color], alpha), width=2)
    return _composite(img, layer, 3)


def add_vignette(img, *, strength=90):
    layer = Image.new("RGBA", (PX_W, PX_H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(layer)
    max_r = int(math.hypot(PX_W, PX_H)) // 2
    cx, cy = PX_W // 2, PX_H // 2
    for r in range(max_r, max_r // 3, -2):
        t = (r - max_r // 3) / (max_r - max_r // 3)
        a = int(strength * t * t)
        draw.ellipse((cx - r, cy - r, cx + r, cy + r), fill=(0, 0, 0, a))
    return _composite(img, layer, 20)


def add_scanlines(img, *, alpha=5, spacing=4):
    layer, draw = _layer()
    for y in range(0, PX_H, spacing):
        draw.line((0, y, PX_W, y), fill=(255, 255, 255, alpha), width=1)
    return _composite(img, layer)


def add_circuit(img, *, color="aqua", alpha=20, nodes=18):
    layer, draw = _layer()
    c_line = (*COLORS[color], alpha)
    c_node = (*COLORS[color], alpha + 30)
    pts = []
    for i in range(nodes):
        x = 100 + (i * 257 + i * i * 19) % (PX_W - 200)
        y = 100 + (i * 193 + i * i * 13) % (PX_H - 200)
        pts.append((x, y))
        draw.ellipse((x - 5, y - 5, x + 5, y + 5), fill=c_node)
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]
        x2, y2 = pts[i + 1]
        mx = (x1 + x2) // 2
        draw.line((x1, y1, mx, y1), fill=c_line, width=2)
        draw.line((mx, y1, mx, y2), fill=c_line, width=2)
        draw.line((mx, y2, x2, y2), fill=c_line, width=2)
    return _composite(img, layer, 1)


def add_prism(img, *, cx=1400, cy=500, size=120, alpha=22):
    layer, draw = _layer()
    pts = [(cx, cy - size), (cx - size, cy + size // 2), (cx + size, cy + size // 2)]
    draw.polygon(pts, outline=(*COLORS["paper"], alpha + 15), fill=(*COLORS["paper"], 4))
    spectrum = ["coral", "amber", "mint", "aqua", "sky", "violet", "magenta"]
    for i, c in enumerate(spectrum):
        y_off = -size + i * (size * 2 // len(spectrum))
        draw.line((cx + size, cy + y_off, cx + size + 400, cy + y_off - 60 + i * 30), fill=(*COLORS[c], alpha), width=3)
    return _composite(img, layer, 4)


def add_ribbon(img, *, color="royal", alpha=20, y_center=500):
    layer, draw = _layer()
    for w in range(30):
        pts = []
        for x in range(0, PX_W + 1, 12):
            y = y_center + w + int(math.sin(x / 300.0) * 60 + math.sin(x / 150.0) * 20)
            pts.append((x, y))
        draw.line(pts, fill=(*COLORS[color], alpha - w // 2), width=2)
    return _composite(img, layer, 6)


def add_dot_matrix(img, *, color="sky", spacing=48, alpha=12, radius=2):
    layer, draw = _layer()
    c = (*COLORS[color], alpha)
    for x in range(0, PX_W, spacing):
        for y in range(0, PX_H, spacing):
            draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=c)
    return _composite(img, layer)


def cinematic_preset(name: str, accent: str, *, dark: bool) -> Image.Image:
    if name == "light_clean":
        img = gradient(COLORS["paper"], (228, 236, 248))
        img = add_grid(img, 100, 10)
        img = add_glow(img, (120, 200, 760, 840), accent, 34, 90)
        img = add_dot_matrix(img, color=accent, alpha=4, spacing=56, radius=2)
        return img
    if name == "warm_problem":
        img = gradient_radial(COLORS["amber"], COLORS["coral"], cx=1440, cy=360)
        img = add_glow(img, (400, 140, 1600, 860), accent, 30, 65)
        img = add_vignette(img, strength=34)
        return img
    if name == "deep_circuit":
        img = gradient(COLORS["night"], COLORS["navy"])
        img = add_circuit(img, color=accent, alpha=18, nodes=16)
        img = add_glow(img, (1120, 180, 1880, 900), accent, 90, 72)
        img = add_particles(img, color="paper", density=900, alpha_max=12)
        return img
    if name == "tech_grid":
        img = gradient(COLORS["void"], COLORS["deep"])
        img = add_grid(img, 82, 16, color=accent)
        img = add_diagonal_lines(img, 10, color=accent)
        img = add_glow(img, (1220, 160, 1860, 920), accent, 72, 70)
        img = add_particles(img, color="ice", density=820, alpha_max=10)
        return img
    if name == "aurora_dark":
        img = gradient(COLORS["void"], COLORS["navy"])
        img = add_aurora(img, colors=(accent, "sky", "violet"), bands=4, amplitude=60, blur_r=48)
        img = add_glow(img, (880, 120, 1780, 820), accent, 72, 72)
        img = add_particles(img, color="paper", density=1200, alpha_max=12)
        img = add_vignette(img, strength=58)
        return img
    img = gradient(COLORS["night"] if dark else COLORS["paper"], COLORS["deep"] if dark else COLORS["mist"])
    return add_particles(img, color="paper" if dark else accent, density=900, alpha_max=10)


# ════════════════════════════════════════════════════════════════
#  BACKGROUND PRESETS — 28 CINEMATIC BACKGROUNDS
# ════════════════════════════════════════════════════════════════
def save_bg(name, img):
    img.convert("RGB").save(ASSETS / name)


def build_backgrounds():
    ensure_assets()
    configs = []

    # 01 — Official cover: light, professional, elegant
    img = gradient((247, 249, 252), (225, 233, 246))
    img = add_glow(img, (1080, 60, 1840, 860), "electric", 26, 75)
    img = add_grid(img, 108, 12)
    img = add_rings(img, color="royal", alpha=18, center=(1540, 280), start_radius=120, count=5)
    img = add_dot_matrix(img, color="royal", alpha=6, spacing=54)
    configs.append(("01-official.png", img))

    # 02 — Cinematic: dark, dramatic, immersive
    img = cinematic_preset("aurora_dark", "aqua", dark=True)
    img = add_diagonal_lines(img, 16, color="aqua")
    img = add_bokeh(img, colors=("royal", "violet", "sky", "aqua"), count=22, alpha_range=(10, 28))
    configs.append(("02-cinematic.png", img))

    # 03 — Summary: light, clean, mesh
    img = cinematic_preset("light_clean", "royal", dark=False)
    img = add_mesh(img, color="royal", alpha=10, x_offset=-120)
    img = add_topographic(img, color="royal", lines=5, alpha=6)
    configs.append(("03-summary.png", img))

    # 04 — Problem: warm dark, tension
    img = cinematic_preset("warm_problem", "coral", dark=True)
    img = add_constellation(img, color="coral", count=22, alpha=12)
    img = add_particles(img, density=900, alpha_max=10)
    configs.append(("04-problem.png", img))

    # 05 — What is SEN: dark, aqua focus, nodes
    img = gradient(COLORS["night"], COLORS["navy"])
    img = add_glow(img, (980, 90, 1840, 980), "aqua", 90, 62)
    img = add_nodes(img)
    img = add_circuit(img, color="aqua", alpha=12, nodes=12)
    img = add_particles(img, density=1400, alpha_max=12)
    img = add_scanlines(img, alpha=3)
    configs.append(("05-what.png", img))

    # 06 — Scope: light, structured
    img = gradient((246, 248, 252), (232, 238, 248))
    img = add_grid(img, 96, 14)
    img = add_glow(img, (980, 70, 1820, 960), "sky", 26, 68)
    img = add_columns(img, color="electric", alpha=18, start=1080)
    img = add_hexgrid(img, color="royal", alpha=6, size=70)
    configs.append(("06-scope.png", img))

    # 07 — Methodology: dark, mint, rings
    img = cinematic_preset("aurora_dark", "mint", dark=True)
    img = add_rings(img, color="mint", alpha=18, center=(1460, 320), start_radius=110, count=6)
    configs.append(("07-method.png", img))

    # 08 — Stakeholders: light, violet
    img = gradient((247, 249, 253), (233, 239, 248))
    img = add_glow(img, (1030, 50, 1840, 960), "violet", 26, 66)
    img = add_grid(img, 104, 12)
    img = add_waves(img, color="royal", alpha=16, amplitude=18, baseline=890)
    img = add_dot_matrix(img, color="lavender", alpha=7, spacing=60)
    configs.append(("08-stakeholders.png", img))

    # 09 — Weaknesses: dark, dramatic violet
    img = gradient(COLORS["night"], COLORS["deep"])
    img = add_glow(img, (980, 100, 1880, 980), "violet", 95, 68)
    img = add_waves(img, color="aqua", alpha=26, amplitude=26, baseline=800)
    img = add_bokeh(img, colors=("violet", "purple", "lavender"), count=25, alpha_range=(10, 30))
    img = add_particles(img, density=1200, alpha_max=10)
    img = add_vignette(img, strength=50)
    configs.append(("09-weaknesses.png", img))

    # 10 — Model: light, structured flow
    img = cinematic_preset("light_clean", "royal", dark=False)
    img = add_diagonal_lines(img, 12, color="royal")
    img = add_ribbon(img, color="electric", alpha=8, y_center=950)
    configs.append(("10-model.png", img))

    # 11 — Experience: dark, blueprint
    img = gradient(COLORS["night"], COLORS["navy"])
    img = add_glow(img, (980, 100, 1880, 980), "royal", 108, 68)
    img = add_blueprint(img)
    img = add_lens_flare(img, cx=1600, cy=200, alpha=35, size=140)
    img = add_particles(img, density=1400, alpha_max=12)
    configs.append(("11-experience.png", img))

    # 12 — Func auth: light, mesh + hexgrid
    img = gradient((247, 249, 253), (232, 238, 248))
    img = add_grid(img, 90, 16)
    img = add_glow(img, (-100, 360, 760, 1120), "electric", 28, 74)
    img = add_mesh(img, color="royal", alpha=15, x_offset=-240)
    img = add_hexgrid(img, color="electric", alpha=5, size=80)
    configs.append(("12-func-auth.png", img))

    # 13 — Func ops: dark, arcs + circuit
    img = gradient(COLORS["night"], COLORS["deep"])
    img = add_glow(img, (1020, 100, 1860, 980), "aqua", 96, 64)
    img = add_arcs(img, 28, color="sky")
    img = add_circuit(img, color="aqua", alpha=14, nodes=14)
    img = add_particles(img, density=1350, alpha_max=12)
    configs.append(("13-func-ops.png", img))

    # 14 — Func intelligence: light, waves + topographic
    img = gradient((247, 249, 253), (234, 240, 248))
    img = add_glow(img, (980, 70, 1860, 960), "violet", 24, 64)
    img = add_grid(img, 102, 13)
    img = add_waves(img, color="violet", alpha=13, amplitude=14, baseline=890)
    img = add_topographic(img, color="lavender", lines=8, alpha=8)
    configs.append(("14-func-intelligence.png", img))

    # 15 — NFR: dark, mint columns + aurora
    img = gradient(COLORS["night"], (17, 37, 66))
    img = add_glow(img, (980, 90, 1860, 980), "mint", 88, 64)
    img = add_columns(img, color="mint", alpha=18, start=1050)
    img = add_aurora(img, colors=("mint", "emerald", "neon"), bands=3, amplitude=40, blur_r=50)
    img = add_particles(img, density=1500, alpha_max=14)
    configs.append(("15-nfr.png", img))

    # 16 — Feasibility tech: light, blueprint + prism
    img = gradient((246, 248, 252), (232, 239, 249))
    img = add_grid(img, 92, 15)
    img = add_glow(img, (1010, 80, 1860, 960), "electric", 28, 64)
    img = add_blueprint(img, color="royal")
    img = add_prism(img, cx=1620, cy=460, size=80, alpha=12)
    configs.append(("16-feasibility-tech.png", img))

    # 17 — Feasibility business: dark, amber + violet
    img = gradient(COLORS["night"], COLORS["deep"])
    img = add_glow(img, (980, 90, 1880, 980), "amber", 68, 68)
    img = add_glow(img, (1140, 120, 1720, 720), "violet", 56, 52)
    img = add_arcs(img, 26, color="amber", x_shift=80)
    img = add_bokeh(img, colors=("amber", "gold", "rose"), count=20, alpha_range=(10, 28))
    img = add_particles(img, density=1200, alpha_max=10)
    configs.append(("17-feasibility-business.png", img))

    # 18 — Use cases: light, rings
    img = cinematic_preset("light_clean", "sky", dark=False)
    img = add_rings(img, color="sky", alpha=10, center=(1560, 300), start_radius=120, count=5)
    configs.append(("18-usecases.png", img))

    # 19 — ERD: dark, aqua mesh + constellation
    img = gradient(COLORS["night"], COLORS["navy"])
    img = add_glow(img, (980, 80, 1880, 980), "aqua", 96, 62)
    img = add_mesh(img, color="aqua", alpha=18, x_offset=20)
    img = add_constellation(img, color="aqua", count=30, alpha=14)
    img = add_particles(img, density=1400, alpha_max=12)
    configs.append(("19-erd.png", img))

    # 20 — Trust: light, waves
    img = cinematic_preset("light_clean", "violet", dark=False)
    img = add_waves(img, color="violet", alpha=12, amplitude=18, baseline=860)
    img = add_hexgrid(img, color="sky", alpha=4, size=65)
    configs.append(("20-trust.png", img))

    # 21 — Business: dark, violet nodes
    img = gradient(COLORS["night"], COLORS["deep"])
    img = add_glow(img, (1000, 90, 1860, 980), "violet", 94, 64)
    img = add_nodes(img, color="mint")
    img = add_radial_burst(img, cx=1500, cy=350, color="violet", rays=20, alpha=10, length=500)
    img = add_particles(img, density=1350, alpha_max=12)
    configs.append(("21-business.png", img))

    # 22 — Marketing: light, columns + grid
    img = gradient((247, 249, 253), (233, 239, 248))
    img = add_glow(img, (990, 80, 1850, 950), "royal", 26, 64)
    img = add_columns(img, color="royal", alpha=12, start=1020)
    img = add_grid(img, 106, 12)
    img = add_ribbon(img, color="violet", alpha=8, y_center=920)
    configs.append(("22-marketing.png", img))

    # 23 — Advantages: dark, violet arcs + aurora
    img = gradient(COLORS["night"], COLORS["deep"])
    img = add_glow(img, (980, 100, 1880, 980), "aqua", 90, 62)
    img = add_arcs(img, 30, color="violet")
    img = add_aurora(img, colors=("violet", "magenta", "lavender"), bands=3, amplitude=45, blur_r=50)
    img = add_particles(img, density=1300, alpha_max=12)
    configs.append(("23-advantage.png", img))

    # 24 — Roadmap: light, mint mesh
    img = cinematic_preset("light_clean", "emerald", dark=False)
    img = add_mesh(img, color="mint", alpha=12, x_offset=100)
    img = add_topographic(img, color="emerald", lines=5, alpha=6)
    configs.append(("24-roadmap.png", img))

    # 25 — Value: dark, royal blueprint + lens flare
    img = gradient(COLORS["night"], COLORS["navy"])
    img = add_glow(img, (980, 90, 1880, 980), "royal", 104, 68)
    img = add_blueprint(img, color="aqua")
    img = add_lens_flare(img, cx=1650, cy=250, alpha=40, size=160)
    img = add_particles(img, density=1350, alpha_max=10)
    img = add_scanlines(img, alpha=3)
    configs.append(("25-value.png", img))

    # 26 — Defense: light, violet rings
    img = gradient((247, 249, 253), (233, 239, 248))
    img = add_glow(img, (1040, 70, 1840, 960), "violet", 26, 62)
    img = add_grid(img, 108, 12)
    img = add_rings(img, color="violet", alpha=12, center=(1560, 300), start_radius=100, count=6)
    img = add_prism(img, cx=1580, cy=500, size=70, alpha=8)
    configs.append(("26-defense.png", img))

    # 27 — Finale: dark, multi-glow + waves
    img = gradient(COLORS["void"], COLORS["deep"])
    img = add_glow(img, (980, 90, 1880, 980), "amber", 78, 70)
    img = add_glow(img, (1050, 110, 1780, 740), "aqua", 56, 52)
    img = add_waves(img, color="mint", alpha=22, amplitude=24, baseline=840)
    img = add_bokeh(img, colors=("amber", "aqua", "mint", "gold"), count=28, alpha_range=(10, 28))
    img = add_particles(img, density=1500, alpha_max=12)
    img = add_vignette(img, strength=50)
    configs.append(("27-finale.png", img))

    # 28 — Closing: dark, cinematic, full drama
    img = cinematic_preset("aurora_dark", "aqua", dark=True)
    img = add_diagonal_lines(img, 14, color="aqua")
    img = add_lens_flare(img, cx=1550, cy=280, alpha=24, size=110)
    configs.append(("28-closing.png", img))

    # 29 — Hero tension: bold, minimal
    img = gradient_three(COLORS["void"], COLORS["deep"], COLORS["night"])
    img = add_glow(img, (860, 120, 1860, 980), "coral", 70, 90)
    img = add_radial_burst(img, cx=1480, cy=420, color="coral", rays=18, alpha=12, length=520)
    img = add_bokeh(img, colors=("coral", "amber", "rose"), count=18, alpha_range=(8, 22))
    img = add_vignette(img, strength=72)
    configs.append(("29-hero-tension.png", img))

    # 30 — Contrast: editorial split screen
    img = gradient((244, 247, 252), (229, 236, 248))
    img = add_glow(img, (-200, 260, 760, 1120), "royal", 24, 80)
    img = add_glow(img, (1180, 120, 1980, 920), "violet", 20, 72)
    img = add_grid(img, 112, 11)
    img = add_ribbon(img, color="lavender", alpha=8, y_center=180)
    configs.append(("30-contrast.png", img))

    # 31 — Why now: metric impact
    img = cinematic_preset("aurora_dark", "amber", dark=True)
    img = add_rings(img, color="amber", alpha=14, center=(1520, 300), start_radius=100, count=6)
    configs.append(("31-why-now.png", img))

    # 32 — Product focus: premium showcase
    img = gradient(COLORS["void"], COLORS["navy"])
    img = add_glow(img, (980, 90, 1880, 980), "aqua", 104, 68)
    img = add_mesh(img, color="sky", alpha=16, x_offset=80)
    img = add_lens_flare(img, cx=1460, cy=230, alpha=28, size=130)
    img = add_scanlines(img, alpha=3)
    configs.append(("32-product-focus.png", img))

    # 33 — Requirements break: cinematic separator
    img = gradient_three(COLORS["void"], COLORS["deep"], COLORS["night"])
    img = add_glow(img, (1040, 70, 1860, 980), "electric", 92, 66)
    img = add_aurora(img, colors=("aqua", "sky", "electric"), bands=3, amplitude=48, blur_r=46)
    img = add_diagonal_lines(img, 18, color="electric")
    img = add_particles(img, density=1400, alpha_max=12)
    configs.append(("33-requirements-break.png", img))

    # 34 — Diagram spotlight
    img = cinematic_preset("light_clean", "sky", dark=False)
    img = add_hexgrid(img, color="sky", alpha=4, size=74)
    img = add_topographic(img, color="electric", lines=6, alpha=7)
    configs.append(("34-spotlight.png", img))

    total = len(configs)
    for i, (name, image) in enumerate(configs, 1):
        log.info("  🎨 [%d/%d] %s", i, total, name)
        save_bg(name, image)


# ════════════════════════════════════════════════════════════════
#  SLIDE COMPONENTS
# ════════════════════════════════════════════════════════════════
def set_rtl(paragraph):
    ppr = paragraph._p.get_or_add_pPr()
    ppr.set("rtl", "1")


def add_bg(slide, filename):
    slide.shapes.add_picture(str(ASSETS / filename), 0, 0, Inches(SLIDE_W), Inches(SLIDE_H))


def add_text_rtl(slide, left, top, width, height, text, *, font_name=BODY, font_size=20, color_name="paper", bold=False, align=PP_ALIGN.RIGHT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    set_rtl(p)
    r = p.add_run()
    r.text = text
    r.font.name = font_name
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color_name)
    return box


def add_multiline_list(slide, left, top, width, lines, *, color_name="paper", font_size=13, line_gap=0.5):
    for i, line in enumerate(lines):
        add_text_rtl(slide, left, top + i * line_gap, width, 0.32, f"• {line}", font_size=font_size, color_name=color_name)


def add_label(slide, text, left, top, *, color_name="aqua", width=2.05):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("night")
    shape.fill.transparency = 0.2
    shape.line.color.rgb = rgb(color_name)
    shape.line.width = Pt(1.0)
    add_text_rtl(slide, left + 0.14, top + 0.04, width - 0.3, 0.2, text, font_size=9, color_name=color_name, bold=True)


def add_title_cluster(slide, eyebrow, title, subtitle, *, dark=True, label_width=2.1):
    fg = "paper" if dark else "ink"
    muted = "mist" if dark else "grid"
    accent = "aqua" if dark else "royal"
    add_label(slide, eyebrow, 10.22, 0.48, color_name=accent, width=label_width)
    add_text_rtl(slide, 5.6, 0.92, 6.2, 1.42, title, font_name=DISPLAY, font_size=31, color_name=fg, bold=True)
    if subtitle:
        add_text_rtl(slide, 5.64, 2.08, 6.05, 0.76, subtitle, font_size=13, color_name=muted)
    add_accent_bar(slide, 5.62, 2.96, 1.22, color_name=accent, height=0.03)


def add_card(slide, left, top, width, height, *, fill_name="night", border_name="aqua", transparency=0.1, rounded=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shadow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + 0.05), Inches(top + 0.06), Inches(width), Inches(height))
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = rgb("night")
    shadow.fill.transparency = 0.72
    shadow.line.fill.background()
    slide.shapes._spTree.remove(shadow._element)
    slide.shapes._spTree.insert(2, shadow._element)
    card = slide.shapes.add_shape(rounded, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(fill_name)
    card.fill.transparency = transparency
    card.line.color.rgb = rgb(border_name)
    card.line.width = Pt(1.0)
    return card


def add_glass_card(slide, left, top, width, height, *, border_name="aqua", dark=True):
    fill = "night" if dark else "paper"
    trans = 0.22 if dark else 0.08
    return add_card(slide, left, top, width, height, fill_name=fill, border_name=border_name, transparency=trans)


def add_stat_card(slide, left, top, width, height, title, body, *, dark=True, accent="aqua"):
    fill = "night" if dark else "paper"
    fg = "paper" if dark else "ink"
    muted = "mist" if dark else "grid"
    add_card(slide, left, top, width, height, fill_name=fill, border_name=accent, transparency=0.08 if dark else 0.02)
    add_text_rtl(slide, left + 0.2, top + 0.16, width - 0.4, 0.3, title, font_name=DISPLAY, font_size=18, color_name=fg, bold=True)
    add_text_rtl(slide, left + 0.2, top + 0.72, width - 0.4, max(height - 0.9, 0.26), body, font_size=11, color_name=muted)


def add_asset_placeholder(slide, left, top, width, height, path, *, dark=True, accent="aqua", title=None):
    fill = "night" if dark else "paper"
    fg = "paper" if dark else "ink"
    muted = "mist" if dark else "grid"
    label = title or path.stem.replace("-", " ").replace("_", " ").title()
    add_text_rtl(slide, left + 0.18, top + height * 0.28, width - 0.36, 0.32, "Asset Preview", font_name=DISPLAY, font_size=15, color_name=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text_rtl(slide, left + 0.18, top + height * 0.48, width - 0.36, 0.3, label, font_name=DISPLAY, font_size=13, color_name=fg, bold=True, align=PP_ALIGN.CENTER)
    add_text_rtl(slide, left + 0.22, top + height * 0.68, width - 0.44, 0.3, f"غير موجود: {path.name}", font_size=9, color_name=muted, align=PP_ALIGN.CENTER)


def place_image(slide, path, left, top, width, height, *, margin=0.08):
    if not path.exists():
        return
    with Image.open(path) as im:
        sw, sh = im.size
    bw, bh = width - margin * 2, height - margin * 2
    ratio = min(bw / sw, bh / sh)
    dw, dh = sw * ratio, sh * ratio
    x = left + margin + (bw - dw) / 2
    y = top + margin + (bh - dh) / 2
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(dw), Inches(dh))


def add_image_card(slide, path, left, top, width, height, *, border="aqua", fill_name="night", transparency=0.06):
    add_card(slide, left, top, width, height, fill_name=fill_name, border_name=border, transparency=transparency)
    if path.exists():
        place_image(slide, path, left, top, width, height)
        return
    add_asset_placeholder(slide, left, top, width, height, path, dark=(fill_name != "paper"), accent=border)


def add_logo(slide, *, word_left=11.0, word_top=0.72, word_width=0.95, mark_left=None, mark_top=0.82, mark_width=0.54):
    if mark_left is None:
        mark_left = word_left - 0.72
    if LOGO_WORD.exists():
        slide.shapes.add_picture(str(LOGO_WORD), Inches(word_left), Inches(word_top), Inches(word_width), Inches(word_width))
    if LOGO_MARK.exists():
        slide.shapes.add_picture(str(LOGO_MARK), Inches(mark_left), Inches(mark_top), Inches(mark_width), Inches(mark_width))


def add_kpi_chip(slide, left, top, title, value, *, dark=True, accent="aqua"):
    fill = "night" if dark else "paper"
    fg = "paper" if dark else "ink"
    muted = "mist" if dark else "grid"
    add_card(slide, left, top, 2.1, 1.2, fill_name=fill, border_name=accent, transparency=0.08 if dark else 0.02)
    add_text_rtl(slide, left + 0.16, top + 0.14, 1.8, 0.2, title, font_size=10, color_name=muted)
    add_text_rtl(slide, left + 0.16, top + 0.48, 1.7, 0.28, value, font_name=DISPLAY, font_size=18, color_name=fg, bold=True)


def add_timeline_node(slide, left, top, title, body, *, dark=False, accent="royal"):
    fill = "paper" if not dark else "night"
    fg = "ink" if not dark else "paper"
    muted = "grid" if not dark else "mist"
    add_card(slide, left, top, 2.5, 1.8, fill_name=fill, border_name=accent, transparency=0.02 if not dark else 0.08)
    add_text_rtl(slide, left + 0.18, top + 0.18, 2.05, 0.22, title, font_name=DISPLAY, font_size=16, color_name=fg, bold=True)
    add_text_rtl(slide, left + 0.18, top + 0.74, 2.05, 0.72, body, font_size=10, color_name=muted)


def add_pull_quote(slide, left, top, width, quote, *, byline=None, dark=True, accent="aqua"):
    fg = "paper" if dark else "ink"
    muted = "mist" if dark else "grid"
    add_glass_card(slide, left, top, width, 2.15, border_name=accent, dark=dark)
    add_accent_bar(slide, left + width - 0.12, top + 0.18, 0.04, color_name=accent, height=1.75)
    add_text_rtl(slide, left + 0.28, top + 0.28, width - 0.6, 1.1, f"“{quote}”", font_name=DISPLAY, font_size=24, color_name=fg, bold=True)
    if byline:
        add_text_rtl(slide, left + 0.28, top + 1.55, width - 0.6, 0.24, byline, font_size=11, color_name=muted)


def add_ghost_word(slide, text, left, top, width, height, *, color_name="slate", font_size=86, align=PP_ALIGN.LEFT):
    add_text_rtl(slide, left, top, width, height, text, font_name=DISPLAY, font_size=font_size, color_name=color_name, bold=True, align=align)


def add_metric_pillar(slide, left, top, width, height, number, label, body, *, accent="aqua", dark=True):
    fill = "night" if dark else "paper"
    fg = "paper" if dark else "ink"
    muted = "mist" if dark else "grid"
    add_glass_card(slide, left, top, width, height, border_name=accent, dark=dark)
    add_accent_bar(slide, left + width - 0.05, top + 0.14, 0.03, color_name=accent, height=height - 0.28)
    add_text_rtl(slide, left + 0.14, top + 0.16, width - 0.28, 0.42, number, font_name=DISPLAY, font_size=28, color_name=accent, bold=True)
    add_text_rtl(slide, left + 0.14, top + 0.72, width - 0.28, 0.2, label, font_name=DISPLAY, font_size=13, color_name=fg, bold=True)
    add_text_rtl(slide, left + 0.14, top + 1.08, width - 0.28, height - 1.2, body, font_size=10, color_name=muted)


def add_split_panel(slide, left, top, width, height, *, fill_name="paper", border_name="royal", transparency=0.04):
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb(fill_name)
    panel.fill.transparency = transparency
    panel.line.color.rgb = rgb(border_name)
    panel.line.width = Pt(1.6)
    return panel


def add_comparison_lane(slide, left, top, width, title, items, *, accent="royal", dark=False):
    fill = "night" if dark else "paper"
    fg = "paper" if dark else "ink"
    muted = "mist" if dark else "grid"
    add_card(slide, left, top, width, 3.35, fill_name=fill, border_name=accent, transparency=0.08 if dark else 0.02)
    add_text_rtl(slide, left + 0.22, top + 0.18, width - 0.45, 0.24, title, font_name=DISPLAY, font_size=19, color_name=fg, bold=True)
    for idx, item in enumerate(items):
        y = top + 0.78 + idx * 0.48
        add_accent_bar(slide, left + width - 0.18, y + 0.1, 0.04, color_name=accent, height=0.18)
        add_text_rtl(slide, left + 0.22, y, width - 0.6, 0.24, item, font_size=12, color_name=muted if idx else fg)


def add_proof_strip(slide, left, top, width, items, *, dark=True, accent="aqua"):
    fill = "night" if dark else "paper"
    fg = "paper" if dark else "ink"
    muted = "mist" if dark else "grid"
    add_card(slide, left, top, width, 0.86, fill_name=fill, border_name=accent, transparency=0.08 if dark else 0.02)
    seg_w = width / max(len(items), 1)
    for idx, item in enumerate(items):
        x = left + idx * seg_w
        if idx:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(top + 0.15), Inches(x), Inches(top + 0.68))
            line.line.color.rgb = rgb(accent)
            line.line.transparency = 0.35
        add_text_rtl(slide, x + 0.1, top + 0.15, seg_w - 0.2, 0.2, item, font_size=10, color_name=muted if idx else fg, bold=(idx == 0))


def add_spotlight_frame(slide, left, top, width, height, title, *, accent="aqua", dark=True):
    add_glass_card(slide, left, top, width, height, border_name=accent, dark=dark)
    add_label(slide, title, left + width - 2.08, top + 0.12, color_name=accent, width=1.95)
    add_accent_bar(slide, left, top + height - 0.08, width, color_name=accent, height=0.03)


def add_side_annotation(slide, left, top, width, title, body, *, accent="aqua", dark=True):
    fill = "night" if dark else "paper"
    fg = "paper" if dark else "ink"
    muted = "mist" if dark else "grid"
    add_card(slide, left, top, width, 1.16, fill_name=fill, border_name=accent, transparency=0.1 if dark else 0.02)
    add_text_rtl(slide, left + 0.14, top + 0.12, width - 0.28, 0.2, title, font_name=DISPLAY, font_size=13, color_name=fg, bold=True)
    add_text_rtl(slide, left + 0.14, top + 0.46, width - 0.28, 0.46, body, font_size=10, color_name=muted)


def add_orbit_cluster(slide, left, top, size, *, accent="aqua", dark=True):
    tone = "night" if dark else "paper"
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    outer.fill.solid()
    outer.fill.fore_color.rgb = rgb(tone)
    outer.fill.transparency = 0.78 if dark else 0.92
    outer.line.color.rgb = rgb(accent)
    outer.line.width = Pt(1.8)

    mid = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + size * 0.15), Inches(top + size * 0.15), Inches(size * 0.7), Inches(size * 0.7))
    mid.fill.solid()
    mid.fill.fore_color.rgb = rgb("void" if dark else "snow")
    mid.fill.transparency = 0.58 if dark else 0.08
    mid.line.color.rgb = rgb("sky" if accent != "sky" else "violet")
    mid.line.transparency = 0.28
    mid.line.width = Pt(1.2)

    core = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + size * 0.36), Inches(top + size * 0.36), Inches(size * 0.28), Inches(size * 0.28))
    core.fill.solid()
    core.fill.fore_color.rgb = rgb(accent)
    core.fill.transparency = 0.08
    core.line.fill.background()


def add_editorial_band(slide, left, top, width, title, body, *, accent="aqua", dark=True):
    add_glass_card(slide, left, top, width, 1.18, border_name=accent, dark=dark)
    add_accent_bar(slide, left + width - 0.08, top + 0.14, 0.03, color_name=accent, height=0.9)
    add_text_rtl(slide, left + 0.16, top + 0.16, width - 0.36, 0.22, title, font_name=DISPLAY, font_size=15, color_name="paper" if dark else "ink", bold=True)
    add_text_rtl(slide, left + 0.16, top + 0.5, width - 0.36, 0.42, body, font_size=11, color_name="mist" if dark else "grid")


def add_staggered_badges(slide, left, top, badges, *, accent="violet", dark=True):
    fg = "paper" if dark else "ink"
    fill = "night" if dark else "paper"
    for idx, text in enumerate(badges):
        x = left + (idx % 3) * 1.78 + (0.28 if idx % 2 else 0.0)
        y = top + (idx // 3) * 0.55
        add_card(slide, x, y, 1.58, 0.4, fill_name=fill, border_name=accent if idx % 2 == 0 else "sky", transparency=0.08 if dark else 0.02)
        add_text_rtl(slide, x + 0.1, y + 0.08, 1.36, 0.18, text, font_size=9, color_name=fg, bold=True)


def add_radial_kpi_ring(slide, left, top, size, value, label, *, accent="aqua", dark=True):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    outer.fill.solid()
    outer.fill.fore_color.rgb = rgb("night" if dark else "paper")
    outer.fill.transparency = 0.2 if dark else 0.05
    outer.line.color.rgb = rgb(accent)
    outer.line.width = Pt(2.0)
    mid = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left + size * 0.12), Inches(top + size * 0.12), Inches(size * 0.76), Inches(size * 0.76))
    mid.fill.solid()
    mid.fill.fore_color.rgb = rgb("void" if dark else "snow")
    mid.fill.transparency = 0.15 if dark else 0.0
    mid.line.color.rgb = rgb(accent)
    mid.line.transparency = 0.35
    mid.line.width = Pt(1.2)
    add_text_rtl(slide, left + 0.1, top + size * 0.26, size - 0.2, 0.35, value, font_name=DISPLAY, font_size=26, color_name="paper" if dark else "ink", bold=True, align=PP_ALIGN.CENTER)
    add_text_rtl(slide, left + 0.15, top + size * 0.62, size - 0.3, 0.2, label, font_size=10, color_name="mist" if dark else "grid", align=PP_ALIGN.CENTER)


def add_decision_matrix(slide, left, top, width, height, headers, cells, *, accent="royal", dark=False):
    fill = "paper" if not dark else "night"
    fg = "ink" if not dark else "paper"
    muted = "grid" if not dark else "mist"
    add_card(slide, left, top, width, height, fill_name=fill, border_name=accent, transparency=0.02 if not dark else 0.08)
    col_w = width / 2
    row_h = height / 2
    add_accent_bar(slide, left + col_w - 0.01, top + 0.12, 0.02, color_name=accent, height=height - 0.24)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(left + 0.12), Inches(top + row_h), Inches(left + width - 0.12), Inches(top + row_h))
    line.line.color.rgb = rgb(accent)
    line.line.transparency = 0.35
    add_text_rtl(slide, left + 0.18, top + 0.12, col_w - 0.3, 0.18, headers[0], font_name=DISPLAY, font_size=14, color_name=fg, bold=True)
    add_text_rtl(slide, left + col_w + 0.12, top + 0.12, col_w - 0.3, 0.18, headers[1], font_name=DISPLAY, font_size=14, color_name=fg, bold=True)
    add_text_rtl(slide, left + 0.18, top + row_h + 0.12, col_w - 0.3, 0.18, headers[2], font_name=DISPLAY, font_size=14, color_name=fg, bold=True)
    add_text_rtl(slide, left + col_w + 0.12, top + row_h + 0.12, col_w - 0.3, 0.18, headers[3], font_name=DISPLAY, font_size=14, color_name=fg, bold=True)
    add_text_rtl(slide, left + 0.18, top + 0.46, col_w - 0.3, 0.42, cells[0], font_size=10, color_name=muted)
    add_text_rtl(slide, left + col_w + 0.12, top + 0.46, col_w - 0.3, 0.42, cells[1], font_size=10, color_name=muted)
    add_text_rtl(slide, left + 0.18, top + row_h + 0.46, col_w - 0.3, 0.42, cells[2], font_size=10, color_name=muted)
    add_text_rtl(slide, left + col_w + 0.12, top + row_h + 0.46, col_w - 0.3, 0.42, cells[3], font_size=10, color_name=muted)


def add_callout(slide, left, top, width, text, *, accent="aqua", dark=True):
    fill = "night" if dark else "paper"
    fg = "paper" if dark else "ink"
    add_card(slide, left, top, width, 0.5, fill_name=fill, border_name=accent, transparency=0.1 if dark else 0.02)
    add_text_rtl(slide, left + 0.12, top + 0.12, width - 0.24, 0.2, text, font_size=10, color_name=fg, bold=True)


def add_device_frame(slide, path, left, top, width, height, *, accent="aqua"):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    outer.fill.solid()
    outer.fill.fore_color.rgb = rgb("obsidian")
    outer.line.color.rgb = rgb(accent)
    outer.line.width = Pt(1.6)
    screen_left = left + 0.16
    screen_top = top + 0.22
    screen_width = width - 0.32
    screen_height = height - 0.38
    screen = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(screen_left), Inches(screen_top), Inches(screen_width), Inches(screen_height))
    screen.fill.solid()
    screen.fill.fore_color.rgb = rgb("night")
    screen.line.fill.background()
    if path.exists():
        place_image(slide, path, screen_left, screen_top, screen_width, screen_height, margin=0.02)
    else:
        add_asset_placeholder(slide, screen_left, screen_top, screen_width, screen_height, path, dark=True, accent=accent, title="Product Screen")
    notch = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left + width / 2 - 0.45), Inches(top + 0.06), Inches(0.9), Inches(0.1))
    notch.fill.solid()
    notch.fill.fore_color.rgb = rgb("slate")
    notch.line.fill.background()


def add_horizontal_flow(slide, items, *, left=1.0, top=4.2, dark=True):
    colors = ["aqua", "sky", "violet", "amber", "mint", "royal", "electric"]
    cur = left
    for idx, item in enumerate(items):
        add_card(slide, cur, top + (0.55 if idx % 2 else 0.0), 1.86, 0.96, fill_name="night" if dark else "paper", border_name=colors[idx % len(colors)], transparency=0.08 if dark else 0.02)
        add_text_rtl(slide, cur + 0.12, top + 0.18 + (0.55 if idx % 2 else 0.0), 1.52, 0.26, item, font_name=DISPLAY, font_size=12, color_name="paper" if dark else "ink", bold=True)
        if idx < len(items) - 1:
            ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(cur + 1.84), Inches(top + 0.5), Inches(cur + 2.05), Inches(top + 0.5))
            ln.line.color.rgb = rgb("sky" if dark else "royal")
            ln.line.width = Pt(1.3)
        cur += 1.98


def add_accent_bar(slide, left, top, width, *, color_name="aqua", height=0.04):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(color_name)
    bar.line.fill.background()
    return bar


def add_slide_number(slide, num, total=TOTAL_SLIDES, *, dark=True):
    accent = accent_for_slide(num)
    chapter = chapter_for_slide(num)
    mood_dark = dark if dark is not None else mood_for_slide(num) == "dark"
    track = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.52), Inches(7.09), Inches(3.76), Inches(0.08))
    track.fill.solid()
    track.fill.fore_color.rgb = rgb("slate" if mood_dark else "mist")
    track.fill.transparency = 0.55
    track.line.fill.background()
    fill_w = 3.76 * (num / max(total, 1))
    fill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.52), Inches(7.09), Inches(fill_w), Inches(0.08))
    fill.fill.solid()
    fill.fill.fore_color.rgb = rgb(accent)
    fill.line.fill.background()
    add_text_rtl(slide, 6.98, 6.98, 1.25, 0.2, chapter, font_size=8, color_name=accent, align=PP_ALIGN.LEFT)
    add_text_rtl(slide, 12.02, 6.98, 0.7, 0.2, f"{num}/{total}", font_size=8, color_name=accent, align=PP_ALIGN.LEFT)


def add_metric_block(slide, left, top, number, label, *, color_name="aqua", dark=True):
    fg = "paper" if dark else "ink"
    muted = "mist" if dark else "grid"
    add_text_rtl(slide, left, top, 2.0, 0.5, number, font_name=DISPLAY, font_size=36, color_name=color_name, bold=True)
    add_text_rtl(slide, left, top + 0.55, 2.0, 0.22, label, font_size=11, color_name=muted)


def add_footer(slide, *, dark=True):
    c = "ash" if dark else "grid"
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.2), Inches(SLIDE_W), Inches(0.02))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb("slate" if dark else "mist")
    bar.fill.transparency = 0.5
    bar.line.fill.background()
    add_text_rtl(slide, 0.45, 7.12, 4.1, 0.2, "SEN Platform | Arabic Cinematic Deck", font_size=8, color_name=c, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════
#  SLIDES 1–14
# ════════════════════════════════════════════════════════════════
def slide_official_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "01-official.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 1, dark=False)
    add_image_card(slide, UNIV_LOGO, 0.9, 0.75, 2.7, 1.95, border="royal", fill_name="paper", transparency=0.01)
    add_logo(slide, word_left=10.7, word_top=0.68, word_width=1.05, mark_left=9.96, mark_top=0.83, mark_width=0.52)
    add_accent_bar(slide, 4.95, 0.6, 4.4, color_name="royal")
    add_text_rtl(slide, 4.95, 0.7, 4.7, 0.24, "وزارة التعليم العالي والبحث العلمي", font_size=11, color_name="grid", bold=True)
    add_text_rtl(slide, 4.95, 1.02, 4.7, 0.26, "جامعة أزال للتنمية البشرية", font_name=DISPLAY, font_size=18, color_name="ink", bold=True)
    add_text_rtl(slide, 4.95, 1.36, 4.7, 0.22, "كلية الحاسوب وتقنية المعلومات", font_size=11, color_name="grid")
    add_text_rtl(slide, 4.95, 1.64, 4.7, 0.22, "قسم تكنولوجيا المعلومات", font_size=11, color_name="grid")
    add_text_rtl(slide, 3.9, 2.2, 7.25, 1.02, "سين للبناء\nوالخدمات المعمارية", font_name=DISPLAY, font_size=28, color_name="ink", bold=True)
    add_text_rtl(slide, 4.15, 3.46, 6.75, 0.4, "مشروع تخرج لنيل درجة البكالوريوس في تقنية المعلومات", font_size=14, color_name="slate")
    add_card(slide, 0.95, 4.2, 5.45, 2.45, fill_name="paper", border_name="royal", transparency=0.01)
    add_text_rtl(slide, 1.28, 4.42, 4.8, 0.24, "إعداد الطلاب", font_name=DISPLAY, font_size=18, color_name="ink", bold=True)
    add_multiline_list(slide, 1.28, 4.88, 4.6, [
        "إبراهيم علي محمد حيدر الرميش  |  2022110613",
        "سيف الدين أحمد سيف المنصوب  |  2022110221",
        "غيدان عدنان محسن قملان  |  2022110657",
        "أمجد أحمد مهيوب سنان  |  2022110511",
    ], color_name="ink", font_size=12, line_gap=0.42)
    add_card(slide, 6.75, 4.2, 5.45, 2.45, fill_name="paper", border_name="violet", transparency=0.01)
    add_text_rtl(slide, 7.08, 4.42, 4.8, 0.24, "الإشراف والسياق", font_name=DISPLAY, font_size=18, color_name="ink", bold=True)
    add_multiline_list(slide, 7.08, 4.88, 4.7, [
        "إشراف الدكتور: مختار غيلان",
        "الحدود المكانية: العاصمة صنعاء",
        "الحدود الزمانية: 2025 - 2026",
        "عرض نهائي مبني على الوثيقة الأكاديمية والمشروع الفعلي.",
    ], color_name="ink", font_size=12, line_gap=0.42)


def slide_cinematic_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "02-cinematic.png")
    add_footer(slide)
    add_slide_number(slide, 2)
    add_ghost_word(slide, "SEN", 0.55, 0.78, 4.4, 1.4, color_name="slate", font_size=96, align=PP_ALIGN.LEFT)
    add_orbit_cluster(slide, 0.86, 2.04, 3.18, accent="aqua", dark=True)
    add_radial_kpi_ring(slide, 1.55, 2.7, 1.8, "2026", "Graduation", accent="aqua", dark=True)
    if LOGO_MARK.exists():
        slide.shapes.add_picture(str(LOGO_MARK), Inches(9.35), Inches(0.96), Inches(1.15), Inches(1.15))
    if LOGO_WORD.exists():
        slide.shapes.add_picture(str(LOGO_WORD), Inches(10.62), Inches(0.78), Inches(1.1), Inches(1.1))
    add_accent_bar(slide, 6.1, 1.18, 5.05, color_name="aqua", height=0.05)
    add_text_rtl(slide, 5.95, 1.34, 5.35, 1.7, "منصة تعيد\nتنظيم القطاع", font_name=DISPLAY, font_size=38, color_name="paper", bold=True)
    add_text_rtl(slide, 6.02, 3.12, 5.3, 0.7, "بنية رقمية تربط السوق والخدمة والقرار والثقة، وتحوّل الفوضى إلى رحلة تشغيل واضحة.", font_size=16, color_name="mist")
    add_editorial_band(slide, 5.96, 3.95, 5.24, "Thesis Positioning", "ليس متجرًا عامًا؛ بل منصة قطاعية متعددة الأطراف تنطلق من سياق محلي واضح.", accent="violet", dark=True)
    add_metric_pillar(slide, 0.95, 4.95, 2.8, 1.45, "01", "النطاق", "بناء، تشطيب، ديكور، وخدمات.", accent="aqua", dark=True)
    add_metric_pillar(slide, 4.02, 4.95, 2.8, 1.45, "02", "النموذج", "منصة متعددة الأطراف وليست متجرًا فقط.", accent="sky", dark=True)
    add_metric_pillar(slide, 7.09, 4.95, 2.8, 1.45, "03", "الموقع", "قراءة محلية واقعية تنطلق من صنعاء.", accent="violet", dark=True)
    add_staggered_badges(slide, 9.96, 5.08, ["Market", "Trust", "Service", "Product", "Local", "Scalable"], accent="aqua", dark=True)


def slide_executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "03-summary.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 3, dark=False)
    add_title_cluster(slide, "الملخص التنفيذي", "SEN يختصر رحلة معقدة.", "الوثيقة تقدّم المشروع كمنصة رقمية تنظّم طلب مواد البناء والديكور والخدمات المرتبطة بها ضمن سوق محلي يحتاج إلى الوضوح والثقة.", dark=False, label_width=2.35)
    add_stat_card(slide, 0.95, 3.2, 3.35, 2.2, "المشكلة", "تشتت المتاجر والعمال والخدمات، وصعوبة المقارنة والطلب الموثوق داخل السوق التقليدي.", dark=False, accent="royal")
    add_stat_card(slide, 4.06, 3.2, 3.35, 2.2, "الحل", "تطبيق موحّد يتيح تصفح المنتجات، طلب الخدمات، واختيار متجر محدد أو طلب عام يرشّح الأنسب.", dark=False, accent="violet")
    add_stat_card(slide, 7.17, 3.2, 3.35, 2.2, "الأثر", "تقليل الوقت والجهد، رفع الشفافية، وتوسيع وصول المتاجر ومقدمي الخدمات إلى عملاء جدد.", dark=False, accent="sky")
    add_stat_card(slide, 10.28, 3.2, 2.1, 2.2, "النتيجة", "نواة منتج قابلة للنمو إلى منصة مرجعية في القطاع.", dark=False, accent="mint")
    add_accent_bar(slide, 0.95, 5.65, 11.45, color_name="royal", height=0.03)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "04-problem.png")
    add_footer(slide)
    add_slide_number(slide, 4)
    add_title_cluster(slide, "المشكلة", "السوق يعمل… لكن بشكل مشتت.", "الوثيقة تبني المشروع على فجوة حقيقية بين العرض والطلب، وضعف البنية الرقمية، واعتماد مفرط على التواصل التقليدي.", dark=True)
    add_stat_card(slide, 0.95, 3.15, 2.75, 2.25, "تشتت السوق", "المستخدم ينتقل بين متاجر متعددة، عمال، وموردين دون منصة موحدة.", dark=True, accent="coral")
    add_stat_card(slide, 3.85, 3.15, 2.75, 2.25, "ضعف المقارنة", "الأسعار والجودة والتقييمات والقرب لا تظهر في إطار قرار واحد.", dark=True, accent="amber")
    add_stat_card(slide, 6.75, 3.15, 2.75, 2.25, "محدودية الثقة", "الاعتماد على المعرفة الشخصية يضعف التوثيق ويزيد احتمالات التعثر.", dark=True, accent="aqua")
    add_stat_card(slide, 9.65, 3.15, 2.75, 2.25, "فجوة رقمية", "كثير من المتاجر ومقدمي الخدمات لا يملكون قناة رقمية فعالة للوصول للعملاء.", dark=True, accent="violet")


def slide_hero_tension(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "29-hero-tension.png")
    add_footer(slide)
    add_slide_number(slide, 5)
    add_ghost_word(slide, "FRACTURE", 0.6, 0.72, 5.0, 0.8, color_name="slate", font_size=48, align=PP_ALIGN.LEFT)
    add_text_rtl(slide, 1.02, 1.35, 6.2, 1.65, "السوق موجود.\nلكن التنظيم ما يزال غائبًا.", font_name=DISPLAY, font_size=36, color_name="paper", bold=True)
    add_text_rtl(slide, 1.06, 3.3, 5.4, 0.48, "هذه هي الفجوة التي تجعل SEN مرئيًا وقابلًا للحل.", font_size=14, color_name="mist")
    add_metric_pillar(slide, 8.25, 1.6, 1.9, 2.15, "3", "أطراف", "عميل، متجر، مقدم خدمة", accent="amber", dark=True)
    add_metric_pillar(slide, 10.35, 1.6, 1.9, 2.15, "1", "منصة", "رحلة قرار موحدة", accent="aqua", dark=True)
    add_text_rtl(slide, 7.1, 4.55, 5.1, 0.72, "من هنا يبدأ التحول:\nمن فوضى القرار إلى هندسة القرار.", font_name=DISPLAY, font_size=24, color_name="paper", bold=True)
    add_proof_strip(slide, 0.95, 5.8, 11.45, ["بحث متشتت", "ثقة محدودة", "طلب غير منظم", "فرصة رقمية واضحة"], dark=True, accent="coral")


def slide_what_is(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "05-what.png")
    add_logo(slide)
    add_footer(slide)
    add_slide_number(slide, 6)
    add_title_cluster(slide, "تعريف المنصة", "SEN ليس متجرًا فقط.", "هو منصة قطاعية متخصصة تربط المستخدم، والمتجر، ومقدم الخدمة داخل منظومة واحدة لعرض المنتجات والخدمات والطلبات والتقييمات.", dark=True)
    add_stat_card(slide, 0.95, 3.45, 2.8, 1.8, "المستخدم", "يتصفح، يقارن، يطلب، ويتابع.", dark=True, accent="aqua")
    add_stat_card(slide, 3.9, 3.45, 2.8, 1.8, "المتجر", "يعرض المنتجات ويدير الطلبات والمخزون.", dark=True, accent="amber")
    add_stat_card(slide, 6.85, 3.45, 2.8, 1.8, "العامل", "يعرض الخدمات، يقبل الطلبات، ويبني سمعة رقمية.", dark=True, accent="sky")
    add_card(slide, 9.8, 3.08, 2.25, 2.5, fill_name="night", border_name="violet", transparency=0.08, rounded=MSO_AUTO_SHAPE_TYPE.OVAL)
    add_text_rtl(slide, 10.13, 3.7, 1.55, 0.24, "النتيجة", font_name=DISPLAY, font_size=18, color_name="paper", bold=True)
    add_text_rtl(slide, 10.02, 4.24, 1.72, 0.6, "سوق منتجات + خدمات + ثقة + متابعة", font_size=11, color_name="mist")


def slide_before_after(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "30-contrast.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 7, dark=False)
    add_text_rtl(slide, 5.45, 0.95, 6.2, 0.9, "ما الذي يتغير\nعندما يدخل SEN؟", font_name=DISPLAY, font_size=34, color_name="ink", bold=True)
    add_split_panel(slide, 0.9, 2.0, 5.65, 4.6, fill_name="paper", border_name="coral", transparency=0.03)
    add_split_panel(slide, 6.78, 2.0, 5.65, 4.6, fill_name="snow", border_name="aqua", transparency=0.02)
    add_ghost_word(slide, "قبل", 0.98, 2.12, 2.8, 0.7, color_name="ash", font_size=56, align=PP_ALIGN.LEFT)
    add_ghost_word(slide, "بعد", 9.9, 2.12, 2.0, 0.7, color_name="grid", font_size=56, align=PP_ALIGN.RIGHT)
    add_multiline_list(slide, 1.2, 3.0, 4.7, [
        "بحث بين متاجر وعمال بشكل منفصل.",
        "مقارنة ضعيفة بين السعر والجودة والتوفر.",
        "اعتماد مرتفع على العلاقات الشخصية.",
        "متابعة غير موحدة للطلبات والخدمات.",
    ], color_name="ink", font_size=14, line_gap=0.58)
    add_multiline_list(slide, 7.05, 3.0, 4.7, [
        "واجهة موحدة تجمع المتجر والخدمة والمستخدم.",
        "قرار أوضح عبر معلومات منظمة وتقييمات.",
        "طلب مباشر أو طلب عام يرشّح الأنسب.",
        "رحلة تشغيلية قابلة للمتابعة والتقييم.",
    ], color_name="ink", font_size=14, line_gap=0.58)
    add_proof_strip(slide, 3.1, 6.15, 7.15, ["عشوائية", "وضوح", "اتصال منفصل", "منصة مترابطة"], dark=False, accent="violet")


def slide_scope(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "06-scope.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 8, dark=False)
    add_title_cluster(slide, "حدود المشروع", "النطاق محدد… وهذا يعزّز وضوحه.", "الوثيقة تضع للمشروع حدودًا مكانية وزمنية واضحة، وتربطه مباشرة بالسوق المحلي وباحتياجات قطاع البناء والتشطيب والخدمات المعمارية.", dark=False, label_width=2.0)
    add_stat_card(slide, 0.95, 3.18, 2.7, 1.85, "مكانيًا", "تم تطبيق المشروع في العاصمة صنعاء.", dark=False, accent="royal")
    add_stat_card(slide, 3.85, 3.18, 2.7, 1.85, "زمنيًا", "2025 - 2026 كما ورد في الوثيقة.", dark=False, accent="violet")
    add_stat_card(slide, 6.75, 3.18, 2.7, 1.85, "قطاعيًا", "البناء، التشطيب، الديكور، المفروشات، والخدمات المهنية.", dark=False, accent="sky")
    add_stat_card(slide, 9.65, 3.18, 2.7, 1.85, "تشغيليًا", "شراء من متجر محدد أو طلب عام يتيح للنظام اقتراح الأنسب.", dark=False, accent="mint")
    add_card(slide, 0.95, 5.35, 11.45, 1.18, fill_name="paper", border_name="royal", transparency=0.01)
    add_text_rtl(slide, 1.25, 5.72, 10.8, 0.34, "هذا التحديد مهم أكاديميًا لأنه يمنح العرض صدقية، ويحوّل المشروع من فكرة عامة إلى حلّ موجّه لسياق واقعي واضح.", font_size=13, color_name="ink")


def slide_why_now(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "31-why-now.png")
    add_footer(slide)
    add_slide_number(slide, 9)
    add_ghost_word(slide, "NOW", 0.8, 0.72, 3.6, 0.7, color_name="slate", font_size=62, align=PP_ALIGN.LEFT)
    add_text_rtl(slide, 5.95, 1.1, 5.4, 1.15, "التوقيت جزء\nمن قوة المشروع.", font_name=DISPLAY, font_size=35, color_name="paper", bold=True)
    add_text_rtl(slide, 6.02, 2.95, 5.1, 0.48, "SEN لا يظهر في فراغ؛ بل في لحظة يلتقي فيها وجود السوق مع فجوة التنظيم وتأخر الرقمنة.", font_size=14, color_name="mist")
    add_metric_pillar(slide, 0.95, 3.55, 2.55, 1.9, "A", "السوق موجود", "قطاع فعلي ومتكرر الطلب.", accent="amber", dark=True)
    add_metric_pillar(slide, 3.82, 3.55, 2.55, 1.9, "B", "الفجوة موجودة", "السوق يعمل لكنه غير منظم رقميًا.", accent="coral", dark=True)
    add_metric_pillar(slide, 6.69, 3.55, 2.55, 1.9, "C", "الرقمنة متأخرة", "المنصات المتخصصة في هذا القطاع قليلة.", accent="sky", dark=True)
    add_metric_pillar(slide, 9.56, 3.55, 2.55, 1.9, "D", "Sen يأتي الآن", "حل يلتقي مع توقيت مناسب للسوق.", accent="mint", dark=True)
    add_pull_quote(slide, 1.25, 5.82, 11.0, "إذا كانت المشكلة حاضرة والسوق قائمًا والرقمنة متأخرة، فهذه ليست مجرد فكرة.\nإنها لحظة مناسبة لبناء منصة.", dark=True, accent="mint")


def slide_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "07-method.png")
    add_logo(slide)
    add_footer(slide)
    add_slide_number(slide, 10)
    add_title_cluster(slide, "المنهجية", "Agile كانت الاختيار المنطقي.", "الوثيقة تنص صراحة على اعتماد Agile بسبب مرونتها العالية وقدرتها على التكيف مع تغيّر المتطلبات ومشاركة أصحاب المصلحة.", dark=True, label_width=1.85)
    add_image_card(slide, AGILE_DIAGRAM, 0.95, 2.6, 4.25, 3.6, border="mint")
    add_stat_card(slide, 5.55, 3.0, 2.1, 1.8, "تطوير تكراري", "تقسيم العمل إلى دورات قصيرة قابلة للتقييم المستمر.", dark=True, accent="aqua")
    add_stat_card(slide, 7.8, 3.0, 2.1, 1.8, "استجابة للتغيير", "إضافة أو تعديل الخصائص دون تعطيل المسار الكامل.", dark=True, accent="sky")
    add_stat_card(slide, 10.05, 3.0, 2.1, 1.8, "مشاركة أصحاب المصلحة", "تغذية راجعة متواصلة ترفع جودة المنتج النهائي.", dark=True, accent="violet")
    add_glass_card(slide, 5.55, 5.1, 6.6, 1.1, border_name="mint")
    add_text_rtl(slide, 5.86, 5.42, 6.05, 0.32, "بناء تدريجي + تقليل للمخاطر + اختبار متكرر = منهجية مناسبة لمشروع متعدد الأطراف مثل SEN", font_size=13, color_name="mist")


def slide_stakeholders(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "08-stakeholders.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 11, dark=False)
    add_title_cluster(slide, "أصحاب المصلحة", "فهم النظام يبدأ من فهم من يخدمهم.", "الوثيقة توسّع المنصة إلى أربع دوائر تشغيلية: العميل، المتجر، مقدم الخدمة، وفريق التطوير/إدارة النظام.", dark=False, label_width=2.1)
    add_stat_card(slide, 0.95, 3.1, 2.75, 2.2, "العميل", "وصول سريع، مقارنة أوضح، طلب مباشر، ومتابعة مريحة للمنتج أو الخدمة.", dark=False, accent="royal")
    add_stat_card(slide, 3.85, 3.1, 2.75, 2.2, "المتجر", "عرض المنتجات، إدارة المخزون، قبول الطلبات، وتتبع المبيعات.", dark=False, accent="amber")
    add_stat_card(slide, 6.75, 3.1, 2.75, 2.2, "مقدم الخدمة", "عرض الخبرة والتخصص، تنظيم المواعيد، واستقبال الطلبات.", dark=False, accent="violet")
    add_stat_card(slide, 9.65, 3.1, 2.75, 2.2, "الإدارة / الفريق", "إدارة المحتوى، المراقبة، والصلاحيات وضمان جاهزية النظام.", dark=False, accent="sky")
    add_card(slide, 0.95, 5.6, 11.45, 0.92, fill_name="paper", border_name="violet", transparency=0.01)
    add_text_rtl(slide, 1.22, 5.9, 10.9, 0.24, "هذه القراءة أقوى من مجرد عرض أنواع الحسابات، لأنها تشرح القيمة المتبادلة داخل المنظومة الاقتصادية والتشغيلية للمشروع.", font_size=12, color_name="ink")


def slide_weaknesses(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "09-weaknesses.png")
    add_footer(slide)
    add_slide_number(slide, 12)
    add_title_cluster(slide, "تحليل الوضع الحالي", "قبل الحل… أين الخلل؟", "الفصل التحليلي في الوثيقة يوضّح نقاط الضعف الحالية للسوق ومسارات العمل التقليدية التي يحاول المشروع إعادة تنظيمها.", dark=True, label_width=2.2)
    add_stat_card(slide, 0.95, 3.1, 2.75, 2.25, "الوقت والجهد", "الطرق التقليدية تستهلك وقتًا كبيرًا على العميل والمتجر ومقدم الخدمة.", dark=True, accent="amber")
    add_stat_card(slide, 3.85, 3.1, 2.75, 2.25, "شفافية محدودة", "الوصول إلى معلومات دقيقة عن الجودة والسعر والتوفر ليس سهلًا.", dark=True, accent="aqua")
    add_stat_card(slide, 6.75, 3.1, 2.75, 2.25, "تجزئة السوق", "لا توجد منصة موحدة للمقارنة أو الطلب السريع من أكثر من مصدر.", dark=True, accent="violet")
    add_stat_card(slide, 9.65, 3.1, 2.75, 2.25, "إدارة ضعيفة", "متابعة الطلبات والمخزون والمواعيد ما تزال غير متكاملة.", dark=True, accent="sky")


def slide_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "10-model.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 13, dark=False)
    add_title_cluster(slide, "النموذج التشغيلي", "كيف يتحرك المستخدم داخل SEN؟", "النظام لا يقدّم واجهة عشوائية؛ بل رحلة منظمة تبدأ بالدور وتنتهي بالتقييم والثقة المتراكمة.", dark=False, label_width=2.05)
    add_horizontal_flow(slide, ["إنشاء حساب", "اختيار الدور", "استكشاف", "مقارنة", "طلب", "متابعة", "تقييم"], left=0.96, top=4.15, dark=False)
    add_card(slide, 0.95, 2.95, 3.55, 0.9, fill_name="paper", border_name="royal", transparency=0.01)
    add_text_rtl(slide, 1.18, 3.24, 3.1, 0.28, "شراء من متجر محدد", font_name=DISPLAY, font_size=15, color_name="ink", bold=True)
    add_card(slide, 4.9, 2.95, 3.55, 0.9, fill_name="paper", border_name="violet", transparency=0.01)
    add_text_rtl(slide, 5.13, 3.24, 3.1, 0.28, "أو طلب عام يرشّح الأنسب", font_name=DISPLAY, font_size=15, color_name="ink", bold=True)
    add_card(slide, 8.85, 2.95, 3.55, 0.9, fill_name="paper", border_name="sky", transparency=0.01)
    add_text_rtl(slide, 9.08, 3.24, 3.1, 0.28, "ثم سجل + تقييم + ثقة", font_name=DISPLAY, font_size=15, color_name="ink", bold=True)


def slide_experience_focus(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "32-product-focus.png")
    add_logo(slide)
    add_footer(slide)
    add_slide_number(slide, 14)
    add_ghost_word(slide, "PRODUCT", 0.8, 0.75, 5.2, 0.6, color_name="slate", font_size=42, align=PP_ALIGN.LEFT)
    add_text_rtl(slide, 6.1, 0.98, 5.4, 0.9, "شاشة واحدة…\nرسالة أوضح.", font_name=DISPLAY, font_size=36, color_name="paper", bold=True)
    add_text_rtl(slide, 6.15, 2.15, 4.95, 0.42, "نركز على واجهة واحدة ونشرح لماذا تصلح كسطح قرار لا كسطح عرض فقط.", font_size=14, color_name="mist")
    add_device_frame(slide, SHOT_HOME, 0.8, 2.0, 6.55, 4.9, accent="aqua")
    add_metric_pillar(slide, 7.55, 2.55, 2.0, 1.55, "01", "اكتشاف", "واجهة رئيسية تدمج التصفح والاستكشاف والعروض.", accent="aqua", dark=True)
    add_metric_pillar(slide, 9.72, 2.55, 2.0, 1.55, "02", "قرار", "تركيز بصري على الفئات والعناصر القابلة للقرار.", accent="sky", dark=True)
    add_metric_pillar(slide, 7.55, 4.45, 2.0, 1.55, "03", "انتقال", "إخراج يسهّل الانتقال من الاكتشاف إلى الطلب.", accent="mint", dark=True)
    add_side_annotation(slide, 9.72, 4.45, 2.0, "UI Focus", "يوجه العين إلى الخطوة التالية بوضوح.", accent="violet", dark=True)


def slide_experience(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "11-experience.png")
    add_logo(slide)
    add_footer(slide)
    add_slide_number(slide, 15)
    add_title_cluster(slide, "المنتج الفعلي", "المشروع يمتلك سطحًا بصريًا حقيقيًا.", "هذه اللقطات مأخوذة من التطبيق نفسه، وتحوّل العرض من خطاب نظري إلى منتج يمكن استعراضه ومناقشته.", dark=True, label_width=2.05)
    add_device_frame(slide, SHOT_HOME, 4.55, 2.15, 3.2, 4.82, accent="aqua")
    add_image_card(slide, SHOT_SECTIONS, 0.95, 2.45, 3.15, 1.92, border="royal")
    add_image_card(slide, SHOT_LOGIN, 8.25, 2.45, 3.15, 1.92, border="sky")
    add_image_card(slide, SHOT_ACCOUNT, 0.95, 4.72, 3.15, 1.92, border="violet")
    add_editorial_band(slide, 8.25, 4.72, 4.15, "Decision Surface", "الواجهة الأقوى ليست الأجمل فقط، بل التي تقود المستخدم بسرعة من الاكتشاف إلى القرار ثم المتابعة.", accent="mint", dark=True)
    add_side_annotation(slide, 4.18, 2.85, 0.9, "01", "Home", accent="aqua", dark=True)
    add_side_annotation(slide, 7.9, 2.85, 0.9, "02", "Flow", accent="sky", dark=True)
    add_text_rtl(slide, 1.08, 4.48, 2.7, 0.2, "الأقسام", font_name=DISPLAY, font_size=13, color_name="paper", bold=True)
    add_text_rtl(slide, 9.0, 4.48, 2.7, 0.2, "الدخول", font_name=DISPLAY, font_size=13, color_name="paper", bold=True)
    add_text_rtl(slide, 1.08, 6.76, 2.7, 0.2, "نوع الحساب", font_name=DISPLAY, font_size=13, color_name="paper", bold=True)


def slide_requirements_break(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "33-requirements-break.png")
    add_footer(slide)
    add_slide_number(slide, 16)
    add_ghost_word(slide, "SYSTEM", 0.72, 0.92, 4.8, 0.7, color_name="slate", font_size=54, align=PP_ALIGN.LEFT)
    add_text_rtl(slide, 6.0, 1.45, 5.8, 1.3, "من الفكرة\nإلى البنية", font_name=DISPLAY, font_size=40, color_name="paper", bold=True)
    add_text_rtl(slide, 6.02, 4.12, 5.5, 0.82, "الآن ننتقل من شرح القيمة إلى إثباتها: المتطلبات، السلوك، والنظام الذي يجعل SEN قابلاً للتنفيذ.", font_size=16, color_name="mist")
    add_proof_strip(slide, 1.0, 5.8, 11.35, ["Roles", "Orders", "Search", "Support", "Trust"], dark=True, accent="electric")


def slide_functional_auth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "12-func-auth.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 17, dark=False)
    add_title_cluster(slide, "المتطلبات الوظيفية", "المصادقة، الأدوار، والخدمات", "الوثيقة لا تكتفي بوصف عام؛ بل تحدد وظائف واضحة تمثل العمود الفقري للنظام.", dark=False, label_width=2.3)
    add_image_card(slide, AUTH_SEQUENCE, 0.95, 2.2, 5.15, 4.9, border="royal", fill_name="paper", transparency=0.0)
    add_stat_card(slide, 6.45, 2.45, 2.75, 1.9, "تسجيل ومصادقة", "إنشاء حساب، تسجيل دخول، إعادة تعيين كلمة المرور، وتحديد صلاحيات حسب نوع الحساب.", dark=False, accent="royal")
    add_stat_card(slide, 9.45, 2.45, 2.75, 1.9, "إدارة الخدمات", "تصفح العمال والخدمات، قبول/رفض الطلب، منشورات وعروض، وتقييم بعد التنفيذ.", dark=False, accent="violet")
    add_stat_card(slide, 6.45, 4.75, 5.75, 1.95, "القيمة هنا", "هذا المخطط من الوثيقة يوضح أن المصادقة ليست شاشة دخول فقط؛ بل بوابة توزيع الأدوار وفتح الواجهات المناسبة لكل فئة داخل النظام.", dark=False, accent="sky")


def slide_functional_ops(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "13-func-ops.png")
    add_logo(slide)
    add_footer(slide)
    add_slide_number(slide, 18)
    add_title_cluster(slide, "المتطلبات الوظيفية", "المنتجات، المتاجر، الطلبات، والدفع", "هذه الطبقة تنقل SEN من فكرة عرض إلى نظام تشغيلي حقيقي يربط المخزون والطلب والحالة.", dark=True, label_width=2.3)
    add_stat_card(slide, 0.95, 3.0, 2.85, 2.15, "إدارة المنتجات", "إضافة، تعديل، حذف، تفاصيل كاملة، تحديث الكميات، وبحث وتصفية.", dark=True, accent="aqua")
    add_stat_card(slide, 3.95, 3.0, 2.85, 2.15, "إدارة المتاجر", "إنشاء متجر، تحديث بياناته، مراجعة من الإدارة، وتصفح حسب الفئة أو الموقع.", dark=True, accent="sky")
    add_stat_card(slide, 6.95, 3.0, 2.85, 2.15, "الطلبات والدفع", "طلب من متجر محدد، حالات تشغيلية واضحة، وسجل دفع مرتبط برقم الطلب.", dark=True, accent="amber")
    add_stat_card(slide, 9.95, 3.0, 2.45, 2.15, "خيارات الدفع", "الدفع عند الاستلام الآن، مع محافظ إلكترونية مستقبلًا مثل كريمي وون كاش.", dark=True, accent="violet")
    add_glass_card(slide, 0.95, 5.55, 11.45, 0.9, border_name="aqua")
    add_text_rtl(slide, 1.22, 5.86, 10.9, 0.24, "الوثيقة هنا تقرأ المنصة كنظام تجارة قطاعي متكامل، لا مجرد واجهة لعرض عناصر ثابتة.", font_size=12, color_name="mist")


def slide_functional_intelligence(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "14-func-intelligence.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 19, dark=False)
    add_title_cluster(slide, "المتطلبات الوظيفية", "البحث، التواصل، التخصيص، والدعم", "هذا الجزء يوسّع المنتج من كونه marketplace إلى منصة أكثر ذكاءً وتفاعلاً واحتواءً للاستخدام اليومي.", dark=False, label_width=2.3)
    add_stat_card(slide, 0.95, 3.0, 2.75, 2.15, "ميزات المشتري", "عربة تسوق، مفضلة، سجل طلبات، ورسائل تأكيد.", dark=False, accent="royal")
    add_stat_card(slide, 3.85, 3.0, 2.75, 2.15, "التواصل", "مراسلة داخل التطبيق، إشعارات، وتحديثات لحالة الطلب والعروض.", dark=False, accent="violet")
    add_stat_card(slide, 6.75, 3.0, 2.75, 2.15, "البحث الذكي", "فلاتر، اقتراحات فورية، سجل بحث، بحث باسم العامل، وفرز بالمعايير المختلفة.", dark=False, accent="sky")
    add_stat_card(slide, 9.65, 3.0, 2.75, 2.15, "التخصيص والدعم", "ثيم فاتح/داكن، توصيات مخصصة، ملاحظات، FAQ، وقنوات دعم مباشرة.", dark=False, accent="mint")
    add_card(slide, 1.45, 5.55, 10.45, 0.84, fill_name="paper", border_name="royal", transparency=0.01)
    add_text_rtl(slide, 1.72, 5.84, 9.9, 0.24, "هذه الطبقة مهمة لأنها تنقل SEN من مشروع واجهات إلى منتج قابل للاستخدام المتكرر والنمو في الثقة والاحتفاظ بالمستخدمين.", font_size=12, color_name="ink")


# ════════════════════════════════════════════════════════════════
#  SLIDES 15–28
# ════════════════════════════════════════════════════════════════
def slide_nfr(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "15-nfr.png")
    add_logo(slide)
    add_footer(slide)
    add_slide_number(slide, 20)
    add_title_cluster(slide, "المتطلبات غير الوظيفية", "ما الذي يجعل النظام صالحًا للاستخدام فعليًا؟", "الأداء والأمان والقابلية للصيانة ليست تفاصيل ثانوية؛ بل شروط بقاء لأي منصة تتعامل مع معاملات حقيقية.", dark=True, label_width=2.5)
    add_stat_card(slide, 0.95, 3.05, 2.75, 2.15, "سهولة الاستخدام", "واجهة بسيطة، واضحة، ومتوافقة مع مختلف الأجهزة والمنصات.", dark=True, accent="mint")
    add_stat_card(slide, 3.85, 3.05, 2.75, 2.15, "السرعة والاستجابة", "تحميل سريع للبيانات واستجابة فورية لإجراءات المستخدم.", dark=True, accent="aqua")
    add_stat_card(slide, 6.75, 3.05, 2.75, 2.15, "الأمان والخصوصية", "تشفير، مصادقة، وصلاحيات وصول تحمي البيانات الشخصية والتشغيلية.", dark=True, accent="sky")
    add_stat_card(slide, 9.65, 3.05, 2.75, 2.15, "الثبات والصيانة", "استقرار تشغيلي، دعم مستمر، وقدرة على التحديث دون تعطيل طويل للخدمة.", dark=True, accent="violet")


def slide_feasibility_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "16-feasibility-tech.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 21, dark=False)
    add_title_cluster(slide, "الجدوى", "الجدوى الفنية واضحة ومقنعة.", "الوثيقة تعرض المشروع كحل قابل للتنفيذ تقنيًا مع مزيج من تطبيق جوال، قاعدة بيانات، ولوحة إدارة ويب.", dark=False, label_width=1.45)
    add_stat_card(slide, 0.95, 3.0, 2.75, 2.15, "التقنيات", "Flutter للتطبيق، MySQL للبيانات، وLaravel للوحة التحكم كما ورد في الوثيقة.", dark=False, accent="royal")
    add_stat_card(slide, 3.85, 3.0, 2.75, 2.15, "الأمان", "إدارة صلاحيات، حماية للبيانات، وفصل بين وصول العميل والمتجر ومقدم الخدمة.", dark=False, accent="violet")
    add_stat_card(slide, 6.75, 3.0, 2.75, 2.15, "القابلية للتوسع", "إضافة منتجات وخدمات جديدة والتعامل مع عدد أكبر من المستخدمين دون كسر البنية.", dark=False, accent="sky")
    add_stat_card(slide, 9.65, 3.0, 2.75, 2.15, "بيئة التطوير", "VS Code وAndroid Studio وXAMPP ضمن بيئة تطوير محلية متوقعة للمشروع.", dark=False, accent="mint")
    add_image_card(slide, TIMELINE_CHART, 8.4, 5.42, 3.95, 1.35, border="royal", fill_name="paper", transparency=0.0)
    add_text_rtl(slide, 0.95, 5.55, 6.7, 0.78, "حتى مع وجود طموح مستقبلي أكبر من النسخة الحالية، فإن الوثيقة تثبت أن المشروع مبني على خيارات تقنية مفهومة وقابلة للتنفيذ، لا على تصور نظري مبهم.", font_size=12, color_name="ink")


def slide_feasibility_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "17-feasibility-business.png")
    add_footer(slide)
    add_slide_number(slide, 22)
    add_title_cluster(slide, "الجدوى", "اقتصاديًا وتشغيليًا وزمنيًا… المشروع منطقي.", "قوة الوثيقة أنها لا تتحدث عن الفكرة فقط، بل عن قابلية التشغيل والالتزام الزمني والفرصة السوقية.", dark=True, label_width=1.45)
    add_stat_card(slide, 0.95, 3.0, 2.75, 2.2, "اقتصاديًا", "تكلفة أولية ممكنة، مع قابلية لتحويل المنصة إلى نموذج دخل عبر العمولات والاشتراكات والإعلانات.", dark=True, accent="amber")
    add_stat_card(slide, 3.85, 3.0, 2.75, 2.2, "تشغيليًا", "واجهة سهلة، دعم فني، لوحة مركزية، وعمليات أكثر سرعة وكفاءة لجميع الأطراف.", dark=True, accent="aqua")
    add_stat_card(slide, 6.75, 3.0, 2.75, 2.2, "زمنيًا", "تحليل، تصميم، برمجة، اختبار، وإطلاق ضمن إطار زمني واضح قابل للمتابعة.", dark=True, accent="sky")
    add_stat_card(slide, 9.65, 3.0, 2.75, 2.2, "تسويقيًا", "فجوة حقيقية في السوق المحلي مع طلب متزايد على المنصات الرقمية في هذا القطاع.", dark=True, accent="violet")


def slide_usecases(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "18-usecases.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 23, dark=False)
    add_title_cluster(slide, "النماذج والمخططات", "Use Case يثبت منطق النظام.", "الوثيقة تحتوي مخططات مستقلة للعميل، ومقدم الخدمة، وصاحب المتجر. وهذه المخططات تجعل سلوك المنصة قابلاً للشرح والدفاع الأكاديمي.", dark=False, label_width=2.55)
    add_image_card(slide, USECASE_USER, 0.95, 2.65, 3.8, 3.95, border="royal", fill_name="paper", transparency=0.0)
    add_image_card(slide, USECASE_WORKER, 4.77, 2.65, 3.8, 3.95, border="violet", fill_name="paper", transparency=0.0)
    add_image_card(slide, USECASE_STORE, 8.59, 2.65, 3.8, 3.95, border="amber", fill_name="paper", transparency=0.0)
    add_text_rtl(slide, 1.32, 6.65, 2.9, 0.2, "حالة استخدام العميل", font_name=DISPLAY, font_size=14, color_name="ink", bold=True)
    add_text_rtl(slide, 5.1, 6.65, 3.1, 0.2, "حالة استخدام مقدم الخدمة", font_name=DISPLAY, font_size=14, color_name="ink", bold=True)
    add_text_rtl(slide, 9.1, 6.65, 2.8, 0.2, "حالة استخدام المتجر", font_name=DISPLAY, font_size=14, color_name="ink", bold=True)


def slide_diagram_spotlight(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "34-spotlight.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 24, dark=False)
    add_text_rtl(slide, 5.6, 0.95, 6.0, 0.95, "ما الذي يهم فعلًا\nداخل المخططات؟", font_name=DISPLAY, font_size=33, color_name="ink", bold=True)
    add_text_rtl(slide, 5.7, 2.0, 5.4, 0.38, "لسنا هنا لعرض المخطط فقط، بل لاستخراج ما يثبته عن منطق النظام.", font_size=13, color_name="grid")
    user_focus = crop_asset(USECASE_USER, (150, 20, 618, 640), "usecase-user-focus")
    erd_focus = crop_asset(ERD_DIAGRAM, (300, 80, 910, 520), "erd-focus-core")
    add_spotlight_frame(slide, 0.95, 2.1, 5.45, 4.6, "Zoom 01", accent="royal", dark=False)
    place_image(slide, user_focus, 1.07, 2.26, 5.18, 3.35, margin=0.02)
    add_callout(slide, 1.2, 5.72, 4.95, "المهم هنا: الرحلة لا تتوقف عند التصفح؛ بل تمتد إلى السلة والطلب والتواصل والإشعارات.", accent="royal", dark=False)
    add_spotlight_frame(slide, 6.7, 2.1, 5.65, 4.6, "Zoom 02", accent="aqua", dark=False)
    place_image(slide, erd_focus, 6.82, 2.26, 5.38, 3.35, margin=0.02)
    add_callout(slide, 6.98, 5.72, 5.0, "ما يهم هنا: Users, Orders, Products, Stores, Workers ليست كيانات منفصلة؛ بل شبكة تشغيل واحدة.", accent="aqua", dark=False)


def slide_erd(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "19-erd.png")
    add_logo(slide)
    add_footer(slide)
    add_slide_number(slide, 25)
    add_title_cluster(slide, "النمذجة البيانية", "ERD يكشف عمق النظام.", "مخطط العلاقات بين الكيانات في الوثيقة يوضح أن المشروع يفكر في البيانات كوحدات مترابطة: مستخدمون، متاجر، منتجات، طلبات، عمال، إشعارات، تقييمات، وعروض.", dark=True, label_width=2.25)
    add_image_card(slide, ERD_DIAGRAM, 0.95, 2.55, 7.55, 4.75, border="aqua", fill_name="night", transparency=0.02)
    add_stat_card(slide, 8.85, 2.8, 3.35, 1.55, "قاعدة البيانات", "MySQL كما ورد في الوثيقة مع مفاتيح أساسية وأجنبية ومؤشرات وقيود تحقق.", dark=True, accent="sky")
    add_stat_card(slide, 8.85, 4.55, 3.35, 1.55, "ما الذي يعنيه هذا؟", "المنصة ليست واجهات فقط؛ بل تمتلك تصورًا واضحًا لتدفق البيانات وتكامل الكيانات.", dark=True, accent="mint")
    add_stat_card(slide, 8.85, 6.3, 3.35, 0.88, "القيمة الدفاعية", "وجود ERD يرفع جودة النقاش الأكاديمي ويقوّي الحديث عن التوسع المستقبلي.", dark=True, accent="violet")


def slide_trust(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "20-trust.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 26, dark=False)
    add_title_cluster(slide, "الثقة وضبط الجودة", "المنصة تحتاج إلى ثقة قابلة للقياس.", "المشروع يعالج الثقة عبر أكثر من طبقة: ملفات منظمة، تقييمات، حالات طلب، إشعارات، ودعم.", dark=False, label_width=2.4)
    add_stat_card(slide, 0.95, 3.12, 2.75, 2.12, "ملفات تعريف واضحة", "بيانات منظمة للعميل والمتجر ومقدم الخدمة تبني هوية رقمية قابلة للتحقق.", dark=False, accent="royal")
    add_stat_card(slide, 3.85, 3.12, 2.75, 2.12, "حالات تشغيلية", "جديد، قيد التنفيذ، مكتمل، ملغي… وضوح في دورة حياة الطلب.", dark=False, accent="violet")
    add_stat_card(slide, 6.75, 3.12, 2.75, 2.12, "مراجعات وتقييمات", "تغذية راجعة موثقة تعزز القرار وترفع الجودة مع الوقت.", dark=False, accent="sky")
    add_stat_card(slide, 9.65, 3.12, 2.75, 2.12, "دعم وملاحظات", "FAQ، تواصل دعم، وملاحظات مستخدمين لتحسين التجربة واستدامتها.", dark=False, accent="mint")


def slide_business_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "21-business.png")
    add_logo(slide)
    add_footer(slide)
    add_slide_number(slide, 27)
    add_title_cluster(slide, "نموذج الأعمال", "SEN قابل لأن يصبح شركة، لا مشروعًا فقط.", "الوثيقة توضح أكثر من مسار دخل، وهذا يجعل العرض أقوى أمام المستثمرين واللجان معًا.", dark=True, label_width=2.05)
    add_stat_card(slide, 0.95, 3.1, 2.75, 2.15, "عمولات", "نسبة على الطلبات المكتملة بين المستخدم والمتجر أو مقدم الخدمة.", dark=True, accent="aqua")
    add_stat_card(slide, 3.85, 3.1, 2.75, 2.15, "اشتراكات", "خطط للمتاجر المميزة أو لمقدمي الخدمات الاحترافيين.", dark=True, accent="sky")
    add_stat_card(slide, 6.75, 3.1, 2.75, 2.15, "إعلانات وإبراز", "ظهور مدفوع داخل نتائج البحث أو الأقسام والعروض الترويجية.", dark=True, accent="amber")
    add_stat_card(slide, 9.65, 3.1, 2.75, 2.15, "خدمات تشغيلية", "رسوم إضافية للوجستيات أو الدعم أو المزايا المتقدمة مستقبلًا.", dark=True, accent="violet")


def slide_marketing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "22-marketing.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 28, dark=False)
    add_title_cluster(slide, "الاستراتيجية التسويقية", "النمو لن يأتي بالصدفة.", "الفصل التسويقي في الوثيقة يقدم مسارات واضحة لاكتساب السوق وبناء الكتلة الحرجة الأولى للمنصة.", dark=False, label_width=2.45)
    add_stat_card(slide, 0.95, 3.15, 2.75, 2.1, "التسويق الرقمي", "إعلانات موجّهة عبر وسائل التواصل الاجتماعي للمستخدمين المهتمين بالبناء والديكور والخدمات.", dark=False, accent="royal")
    add_stat_card(slide, 3.85, 3.15, 2.75, 2.1, "SEO", "الظهور في نتائج البحث عند البحث عن المنتجات والخدمات ذات الصلة.", dark=False, accent="violet")
    add_stat_card(slide, 6.75, 3.15, 2.75, 2.1, "تواصل مباشر", "استقطاب المتاجر ومقدمي الخدمات عبر التواصل المباشر والشراكات الأولية.", dark=False, accent="sky")
    add_stat_card(slide, 9.65, 3.15, 2.75, 2.1, "فترات تجريبية", "تحفيز الانضمام الأولي وبناء العرض داخل المنصة قبل التوسع الأوسع.", dark=False, accent="mint")


def slide_advantages(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "23-advantage.png")
    add_logo(slide)
    add_footer(slide)
    add_slide_number(slide, 29)
    add_title_cluster(slide, "الميزة التنافسية", "لماذا Sen مختلف فعلًا؟", "لأنه لا يعامل السوق كسوق عام؛ بل كسوق قطاعي متخصص يجمع بين المنتج والخدمة والدليل التشغيلي.", dark=True, label_width=2.15)
    add_multiline_list(slide, 6.1, 3.08, 5.7, [
        "تركيز قطاعي واضح على البناء والتشطيب والديكور.",
        "دمج المتجر والعامل والمستخدم في منصة واحدة.",
        "إمكانية الطلب من متجر محدد أو طلب عام مرن.",
        "قابلية نمو جغرافي ووظيفي دون تغيير الفكرة الأساسية.",
        "بناء الثقة عبر المراجعات وحالات الطلب والملفات المنظمة.",
        "قيمة يومية مرتبطة باحتياج حقيقي ومتكرر في السوق.",
    ], color_name="paper", font_size=13, line_gap=0.5)
    if LOGO_MARK.exists():
        slide.shapes.add_picture(str(LOGO_MARK), Inches(1.15), Inches(3.15), Inches(3.35), Inches(3.35))


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "24-roadmap.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 30, dark=False)
    add_title_cluster(slide, "خارطة التطوير", "الوثيقة تفتح باب المرحلة التالية.", "كثير من العناصر في SEN يمكن تحويلها من تصور أكاديمي إلى roadmap منتجية واضحة وقابلة للتنفيذ.", dark=False, label_width=2.1)
    add_timeline_node(slide, 0.95, 3.55, "المرحلة 1", "إكمال المتجر، ربط البيانات الحقيقية، وتثبيت دورة الطلب الأساسية.", accent="royal")
    add_timeline_node(slide, 3.95, 3.05, "المرحلة 2", "خرائط وموقع جغرافي، تواصل مباشر، وجدولة خدمات أكثر دقة.", accent="violet")
    add_timeline_node(slide, 6.95, 3.55, "المرحلة 3", "دفع إلكتروني، توصيات ذكية، ولوحات تحليلية للمتاجر والعمال.", accent="sky")
    add_timeline_node(slide, 9.95, 3.05, "المرحلة 4", "توسع لمدن أخرى، شراكات أكبر، ونموذج تشغيل واستثمار أوسع.", accent="mint")
    for i in range(3):
        x_start = 3.45 + i * 3.0
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x_start), Inches(4.0), Inches(x_start + 0.5), Inches(4.0))
        ln.line.color.rgb = rgb("royal")
        ln.line.width = Pt(1.5)


def slide_value(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "25-value.png")
    add_logo(slide)
    add_footer(slide)
    add_slide_number(slide, 31)
    add_title_cluster(slide, "القيمة الكلية", "لماذا المشروع قوي أكاديميًا واستثماريًا واجتماعيًا؟", "SEN يجمع بين حل مشكلة واقعية، بنية تقنية قابلة للتوسّع، وفرصة اقتصادية ذات أثر اجتماعي مباشر.", dark=True, label_width=2.35)
    add_stat_card(slide, 0.95, 3.15, 2.75, 2.15, "أكاديميًا", "تحليل متطلبات، نمذجة، جدوى، وهندسة نظام متعدد الأطراف.", dark=True, accent="aqua")
    add_stat_card(slide, 3.85, 3.15, 2.75, 2.15, "اقتصاديًا", "سوق متكرر وقابل للرقمنة ونموذج دخل متعدد القنوات.", dark=True, accent="amber")
    add_stat_card(slide, 6.75, 3.15, 2.75, 2.15, "اجتماعيًا", "تمكين العمال والمتاجر الصغيرة ورفع الشفافية داخل السوق.", dark=True, accent="sky")
    add_stat_card(slide, 9.65, 3.15, 2.75, 2.15, "منتجيًا", "نواة قابلة للتحول من مشروع تخرج إلى منتج ثم إلى شركة.", dark=True, accent="violet")


def slide_defense(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "26-defense.png")
    add_footer(slide, dark=False)
    add_slide_number(slide, 32, dark=False)
    add_title_cluster(slide, "خلاصة الدفاع", "ما الذي ستدافعون عنه أمام اللجنة؟", "أفضل دفاع عن SEN ليس الوصف العام، بل منطق الترابط بين المشكلة والوثيقة والنظام والمخططات والفرصة.", dark=False, label_width=2.05)
    add_card(slide, 0.95, 3.05, 11.45, 3.1, fill_name="paper", border_name="royal", transparency=0.01)
    add_multiline_list(slide, 1.25, 3.42, 10.8, [
        "المشروع يعالج مشكلة محلية حقيقية داخل قطاع واسع ومتكرر الحاجة.",
        "الحل المقترح متخصص، متعدد الأطراف، وأكثر نضجًا من فكرة متجر عادي أو دليل خدمات منفصل.",
        "الوثيقة تدعم المشروع بمنهجية، جدوى، ومتطلبات وظيفية وغير وظيفية ومخططات UML وERD.",
        "المشروع يمتلك سطحًا بصريًا فعليًا من التطبيق نفسه، ما يثبت وجود منتج قابل للعرض لا مجرد فكرة.",
        "النطاق واضح، والرؤية المستقبلية قابلة للتحول إلى roadmap تجارية وتقنية بعد التخرج.",
    ], color_name="ink", font_size=13, line_gap=0.48)


def slide_finale(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "27-finale.png")
    add_logo(slide, word_left=10.85, word_top=0.76, word_width=1.0, mark_left=10.15, mark_top=0.9, mark_width=0.46)
    add_footer(slide)
    add_slide_number(slide, 33)
    add_ghost_word(slide, "SEN", 0.72, 0.86, 3.2, 0.8, color_name="slate", font_size=58, align=PP_ALIGN.LEFT)
    add_orbit_cluster(slide, 0.98, 1.84, 3.0, accent="amber", dark=True)
    add_radial_kpi_ring(slide, 1.62, 2.48, 1.7, "34", "Slides", accent="amber", dark=True)
    add_text_rtl(slide, 5.95, 1.46, 5.8, 1.46, "SEN ليس فكرة\nعابرة للسوق.", font_name=DISPLAY, font_size=38, color_name="paper", bold=True)
    add_editorial_band(slide, 5.96, 3.54, 5.46, "Closing Thesis", "المشروع يربط المشكلة بالسوق، والوثيقة بالنظام، والمنتج بالفرصة. لهذا هو مقنع أكاديميًا ومنتجيًا معًا.", accent="aqua", dark=True)
    add_metric_pillar(slide, 0.95, 5.35, 2.8, 1.45, "01", "القيمة", "حل حقيقي لمشكلة متكررة.", accent="amber", dark=True)
    add_metric_pillar(slide, 4.02, 5.35, 2.8, 1.45, "02", "العمق", "وثيقة، نمذجة، ومنطق دفاعي.", accent="aqua", dark=True)
    add_metric_pillar(slide, 7.09, 5.35, 2.8, 1.45, "03", "المستقبل", "منتج قابل للنمو بعد التخرج.", accent="violet", dark=True)
    add_staggered_badges(slide, 10.0, 5.18, ["Viable", "Local", "Defensible", "Scalable"], accent="amber", dark=True)


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, "28-closing.png")
    add_footer(slide)
    add_slide_number(slide, 34)
    if LOGO_WORD.exists():
        slide.shapes.add_picture(str(LOGO_WORD), Inches(10.72), Inches(0.78), Inches(1.25), Inches(1.25))
    if LOGO_MARK.exists():
        slide.shapes.add_picture(str(LOGO_MARK), Inches(8.55), Inches(1.4), Inches(1.6), Inches(1.6))
    add_text_rtl(slide, 5.98, 1.64, 5.88, 1.46, "شكرًا لكم", font_name=DISPLAY, font_size=35, color_name="paper", bold=True)
    add_pull_quote(slide, 5.96, 3.6, 5.45, "جاهزون للمناقشة، وللانتقال من وثيقة قوية إلى نسخة منتج أكثر نضجًا واتساعًا.", dark=True, accent="aqua")
    add_proof_strip(slide, 0.95, 5.65, 6.2, ["Problem", "Platform", "Proof", "Potential"], dark=True, accent="violet")
    add_accent_bar(slide, 7.9, 6.0, 3.9, color_name="aqua")
    add_text_rtl(slide, 7.9, 6.12, 3.9, 0.24, "SEN Final", font_name=DISPLAY, font_size=20, color_name="aqua", bold=True)


# ════════════════════════════════════════════════════════════════
#  SPEAKER NOTES
# ════════════════════════════════════════════════════════════════
def speaker_notes():
    return [
        "نفتتح العرض بصيغة أكاديمية رسمية كما وردت في الوثيقة نفسها: الجامعة، الكلية، القسم، عنوان المشروع، أسماء الفريق، واسم المشرف. هذه البداية تربط العرض مباشرة بمرجعيته الأكاديمية وتمنحه ثقلًا رسميًا واضحًا.",
        "بعد الغلاف الرسمي ننتقل إلى غلاف سينمائي يعيد تعريف Sen كمنتج، لا كوثيقة فقط. الرسالة المقصودة هنا هي أن المشروع محاولة لإعادة تنظيم قطاع كامل داخل تجربة رقمية واحدة.",
        "في الملخص التنفيذي نضغط جوهر المشروع في ثلاث نقاط: المشكلة، الحل، والأثر. Sen هنا يظهر كمنصة تربط بين المنتجات والخدمات والمتاجر والعمال ضمن سوق محلي يحتاج إلى الوضوح والثقة.",
        "هنا نعرض المشكلة كما هي: السوق يعمل لكنه مشتت. المستخدم ينتقل بين متاجر وعمال ومصادر غير منظمة، ولا يجد إطارًا موحدًا للمقارنة أو التحقق أو المتابعة.",
        "هذه الشريحة المقصود بها خلق التوتر السردي. نختزل المشكلة في عبارة قوية: السوق موجود لكن التنظيم غائب. ومن هذا التوتر تبدأ مبررات المشروع الحقيقية.",
        "في هذه النقطة نعرّف Sen تعريفًا دقيقًا: منصة قطاعية متعددة الأطراف، لا متجرًا فقط ولا دليل خدمات فقط. هي بيئة تجمع المنتج والخدمة والطلب والتقييم داخل إطار واحد.",
        "شريحة قبل وبعد توضّح الفرق بين السوق التقليدي والسوق بعد دخول Sen. هذه المقارنة أهم من أي وصف طويل لأنها تجعل التحول مرئيًا ومباشرًا أمام الجمهور.",
        "تحديد الحدود المكانية والزمانية مهم جدًا أمام اللجنة. الوثيقة تذكر صراحة أن المشروع مطبق في العاصمة صنعاء خلال الفترة 2025 إلى 2026، وهذا يمنح العرض سياقًا واقعيًا موثقًا.",
        "لماذا الآن؟ لأن السوق موجود، والحاجة متكررة، والفجوة الرقمية واضحة، والرقمنة في هذا القطاع ما تزال متأخرة. هذه الشريحة تشرح أن المشروع ليس فقط جيدًا، بل مناسب توقيتًا أيضًا.",
        "الوثيقة اعتمدت Agile، وهذا منطقي لأن المشروع متعدد الأطراف ومتطلباته قابلة للتغير. المنهجية التكرارية تسمح بتقسيم العمل إلى دورات قصيرة مع اختبار مستمر وتغذية راجعة متواصلة.",
        "فهم أصحاب المصلحة يشرح بنية المنصة أفضل من مجرد سرد الحسابات. العميل يريد الوضوح، المتجر يريد قناة بيع منظمة، مقدم الخدمة يريد الوصول والجدولة، والإدارة تريد نظامًا قابلًا للمراقبة.",
        "في تحليل الوضع الحالي، الوثيقة توضح أن الخلل ليس فقط في غياب التطبيق، بل في ضياع الوقت، وضعف الشفافية، وتجزئة السوق، وضعف إدارة الطلبات والمخزون والمواعيد.",
        "هذه الشريحة تعرض النموذج التشغيلي للمنصة. الرحلة تبدأ بالحساب والدور، ثم الاستكشاف، ثم المقارنة، ثم الطلب، ثم المتابعة، ثم التقييم وبناء الثقة.",
        "هنا ننتقل من عرض صور كثيرة إلى التركيز على شاشة واحدة ورسالة واحدة. الهدف هو شرح كيف تقود الواجهة المستخدم بصريًا وسلوكيًا بدل الاكتفاء بجاليري لقطات.",
        "هذه الشريحة تثبت أن المشروع يمتلك سطحًا منتجيًا فعليًا. اللقطات مأخوذة من التطبيق نفسه، ما يعني أننا لا نعرض تصورًا نظريًا فقط، بل واجهات قابلة للمشاهدة والمناقشة.",
        "هذه شريحة فصل بصري بين الرؤية وبين البنية. من هنا نبدأ إثبات أن Sen ليس مجرد فكرة جذابة، بل نظام له متطلبات واضحة وسلوك تشغيلي يمكن الدفاع عنه هندسيًا.",
        "في المتطلبات الوظيفية الأولى نركز على المصادقة وتحديد الأدوار وإدارة الخدمات. مخطط التسجيل والمصادقة من الوثيقة يوضح أن الدخول هو بوابة توزيع المسارات والصلاحيات داخل النظام.",
        "في الطبقة الثانية من المتطلبات ننتقل إلى المنتجات والمتاجر والطلبات والدفع. هنا يتحول المشروع إلى نظام تشغيلي فعلي يتعامل مع المخزون، حالات الطلب، وسجل الدفع.",
        "ثم تأتي طبقة البحث والتواصل والتخصيص والدعم. هذه العناصر تعطي المنصة بعدًا أكثر نضجًا، مثل البحث المتقدم، الرسائل، الإشعارات، التخصيص، والملاحظات والدعم المباشر.",
        "المتطلبات غير الوظيفية تبيّن أننا لا نفكر فقط في ما يفعله النظام، بل في كيف يعمل: السرعة، الأمان، التوافقية، الثبات، وقابلية الصيانة. وهذا عنصر أساسي في أي نقاش هندسي جاد.",
        "الجدوى الفنية في الوثيقة واضحة: Flutter للتطبيق، MySQL للبيانات، Laravel للوحة التحكم، وبيئة تطوير مفهومة. هذا يجعل الحل قابلًا للتنفيذ، لا مجرد فكرة عامة.",
        "ثم نغطي الجدوى الاقتصادية والتشغيلية والزمنية. هذه الطبقات تعزز فكرة أن المشروع ليس تمرينًا برمجيًا، بل مشروعًا له منطق تشغيل وسوق وإطار زمني واقعي.",
        "وجود مخططات Use Case لثلاثة أطراف مختلفة يعطينا قوة كبيرة في الشرح. نستطيع هنا أن نوضح بدقة ماذا يفعل العميل، وماذا يفعل مقدم الخدمة، وماذا يفعل صاحب المتجر.",
        "شريحة الـ spotlight تقرّبنا من المخططات نفسها. بدلاً من عرضها كصور فقط، نوضح ما الذي يهم فيها: ما الذي تثبته عن الرحلة الوظيفية، وما الذي تقوله عن قلب البيانات في النظام.",
        "أما مخطط العلاقات بين الكيانات فيرفع مستوى المشروع تقنيًا. نرى المستخدمين والمتاجر والمنتجات والطلبات والعمال والإشعارات والتقييمات في بنية مترابطة، وهذا مهم جدًا في الدفاع الأكاديمي.",
        "هذه الشريحة تعيد التركيز على الثقة وضبط الجودة. لأن المنصة تتعامل مع معاملات حقيقية، فلا بد من ملفات تعريف واضحة، وحالات تشغيلية، وتقييمات موثقة، وآليات دعم.",
        "نموذج الأعمال يوضح لماذا يمكن النظر إلى Sen كنواة شركة لاحقًا. العمولات، الاشتراكات، الإعلانات، والخدمات التشغيلية تعطي أكثر من مسار للدخل والاستدامة.",
        "الاستراتيجية التسويقية في الوثيقة عملية جدًا: تسويق رقمي، تحسين الظهور في محركات البحث، تواصل مباشر مع المتاجر ومقدمي الخدمات، وفترات تجريبية لجذب العرض الأولي.",
        "الميزة التنافسية الجوهرية هي التخصص القطاعي. Sen لا يحاول أن يكون تطبيقًا عامًا، بل منصة تفهم تفاصيل سوق البناء والتشطيب والخدمات المعمارية وتبني تجربتها على هذا الفهم.",
        "خارطة التطوير هنا تُظهر كيف يمكن نقل المشروع من نسخته الحالية إلى نسخة أكثر نضجًا: إكمال المتجر، ربط بيانات حقيقية، إضافة خرائط وجدولة، ثم دفع إلكتروني وتحليلات، ثم توسع أكبر.",
        "في هذه الشريحة نلخص قيمة المشروع من أربع زوايا: أكاديمية، اقتصادية، اجتماعية، ومنتجية. هذا يعطي العرض اتزانًا مهمًا أمام اللجنة وأمام أي قراءة استثمارية أو سوقية.",
        "هذه هي شريحة الدفاع المختصر. إذا سُئلتم لماذا المشروع قوي، فالإجابة ليست فقط في الواجهات، بل في وضوح المشكلة، ومنطقية الحل، ووجود الوثيقة والمخططات والجدوى والنطاق المحدد.",
        "الفينال يعيد وضع المشروع في إطار أكبر: Sen ليس فكرة عابرة أو شاشة جميلة، بل محاولة حقيقية لبناء طبقة تنظيم رقمية لسوق محلي مشتت وقابل للتطوير والنمو.",
        "نختم بالشكر، مع ترك انطباع واضح أن العرض لم يكتف بشرح الفكرة، بل قدّم مشروعًا موثقًا ومصممًا ومقروءًا من زوايا أكاديمية وتقنية ومنتجية وتسويقية في وقت واحد.",
    ]


# ════════════════════════════════════════════════════════════════
#  COM ENHANCEMENT — transitions, animations, sections, notes
# ════════════════════════════════════════════════════════════════
# Transition effect IDs: 1793=Fade, 3853=Push, 3896=Morph, 3844=Wipe, 3855=Split
TRANSITION_PATTERN = [1793, 3853, 3896, 3844, 3855, 1793, 3896]
# Animation effect IDs: 1=Appear, 10=FadeIn, 23=Float, 53=GrowAndTurn, 42=Swivel
ANIM_FIRST = 23     # Float for first element
ANIM_ODD = 10       # FadeIn for odd elements
ANIM_EVEN = 1       # Appear for even elements

SECTIONS = [
    (1, "الافتتاح الرسمي"), (3, "الملخص"),
    (4, "المشكلة"), (6, "تعريف SEN"),
    (8, "النطاق"), (9, "لماذا الآن؟"),
    (10, "المنهجية"), (11, "أصحاب المصلحة"),
    (12, "التحليل"), (13, "النموذج التشغيلي"),
    (14, "المنتج"), (16, "المتطلبات"),
    (21, "الجدوى"), (23, "المخططات"),
    (26, "الثقة"), (27, "الأعمال"),
    (29, "التميّز"), (30, "التطوير"),
    (32, "الدفاع"), (33, "الختام"),
]


def enhance_with_powerpoint(path: Path) -> None:
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        log.info("  ⏩ pywin32 not available — skipping COM enhancement")
        return

    log.info("  ✨ Enhancing with PowerPoint COM...")
    pythoncom.CoInitialize()
    app = None
    presentation = None
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = 1
        presentation = app.Presentations.Open(str(path), WithWindow=False)
        notes = speaker_notes()
        transitions = [TRANSITION_PATTERN[i % len(TRANSITION_PATTERN)] for i in range(len(notes))]
        slide_specific = {2: 1793, 4: 3844, 5: 3855, 10: 3896, 16: 1793, 23: 3853, 24: 3896, 25: 3896, 33: 1793, 34: 1793}

        try:
            sp = presentation.SectionProperties
            while sp.Count > 0:
                sp.Delete(1, False)
            for si, name in reversed(SECTIONS):
                sp.AddBeforeSlide(si, name)
        except Exception:
            pass

        for si in range(1, presentation.Slides.Count + 1):
            slide = presentation.Slides(si)
            tr = slide.SlideShowTransition
            tr.EntryEffect = slide_specific.get(si, transitions[min(si - 1, len(transitions) - 1)])
            tr.Speed = 3
            tr.AdvanceOnClick = True

            seq = slide.TimeLine.MainSequence
            while seq.Count > 0:
                seq.Item(1).Delete()

            animated = []
            for shi in range(2, slide.Shapes.Count + 1):
                sh = slide.Shapes(shi)
                if sh.Width < 12 or sh.Height < 12:
                    continue
                animated.append(sh)

            for order, shape in enumerate(animated):
                if si in {5, 9, 33}:
                    eid = 53 if order == 0 else (10 if order % 2 else 1)
                elif si in {10, 23, 24, 25}:
                    eid = 23 if order < 2 else (10 if order % 2 else 1)
                else:
                    eid = ANIM_FIRST if order == 0 else (ANIM_ODD if order % 2 else ANIM_EVEN)
                trigger = 2 if order == 0 else 3  # 2=AfterPrevious first, 3=WithPrevious rest
                effect = seq.AddEffect(shape, eid, 0, trigger)
                effect.Timing.Duration = 0.44 if si in {5, 9, 24, 25, 33} else 0.34
                effect.Timing.TriggerDelayTime = min(order * (0.08 if si in {23, 24, 25} else 0.05), 0.58)

            try:
                note_shape = slide.NotesPage.Shapes.Placeholders(2)
                note_shape.TextFrame.TextRange.Text = notes[si - 1]
            except Exception:
                pass

        presentation.Save()
        log.info("  ✅ COM enhancement complete")
    except Exception as e:
        log.warning("  ⚠️ COM enhancement failed: %s", e)
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()


# ════════════════════════════════════════════════════════════════
#  BUILD
# ════════════════════════════════════════════════════════════════
def build() -> None:
    log.info("🚀 SEN Final Deck Builder — Monster Hybrid Edition")
    log.info("━" * 50)

    runtime_dir = ROOT / "_sen_monster_runtime"
    runtime_dir.mkdir(parents=True, exist_ok=True)

    log.info("\n🧠 Building with SlideSpec + Layout Engine...")
    log.info("  • canonical engine: sen_cinematic_engine_monster.py")
    log.info("  • slide count: %d", len(MONSTER_SLIDES))

    deck = MonsterDeck(output_dir=runtime_dir)
    built_path = deck.build(filename=OUTPUT.name)
    final_output = OUTPUT
    try:
        shutil.copyfile(built_path, final_output)
    except PermissionError:
        final_output = ROOT / "SEN Final - Monster Hybrid.pptx"
        shutil.copyfile(built_path, final_output)
        log.warning("  ⚠️ Primary output is locked. Saved fallback copy instead: %s", final_output)
    log.info("\n💾 Saved: %s", final_output)

    enhance_with_powerpoint(final_output)
    log.info("\n🎉 Done! Your hybrid cinematic presentation is ready.")
    log.info("━" * 50)


if __name__ == "__main__":
    build()
